#!/usr/bin/env python3
"""
LG Group PPTX Builder — L-Style Design System
LG 그룹 브랜드 가이드라인을 따르는 프레젠테이션 생성 라이브러리

Usage:
    from lg_pptx_builder import LGPresentation

    prs = LGPresentation()
    prs.add_cover("프로젝트 제목", subtitle="팀명", date="2025.10.23")
    prs.add_toc([("Summary", []), ("시스템 소개", ["항목1", "항목2"])])
    prs.add_content("1.1 시스템 개요", section="Summary", chapter="I. Summary",
                    bullets=["내용1", "내용2"])
    prs.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ─────────────────────────────────────────────
# Design Tokens
# ─────────────────────────────────────────────

class LGColors:
    """LG Group brand color palette"""
    RED = RGBColor(0xA5, 0x00, 0x34)           # Primary - LG RED
    DARK_RED = RGBColor(0xC0, 0x00, 0x00)      # Sub-header badges / labels
    BLACK = RGBColor(0x00, 0x00, 0x00)          # Text primary
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)      # Text secondary
    MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)    # Text tertiary / section labels
    LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)     # Surface / legacy bg
    BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)     # Borders / dividers
    CHARCOAL = RGBColor(0x3C, 0x3C, 0x3C)       # Header bars / timeline
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)           # Background / content box bg
    GREEN = RGBColor(0x2E, 0x7D, 0x32)          # Future / planned items
    ORANGE = RGBColor(0xD4, 0x76, 0x0A)         # Highlight items
    BLUE = RGBColor(0x15, 0x65, 0xC0)           # Accent blue


class HanwhaColors:
    """한화 그룹 브랜드 컬러 팔레트 (Hanwha Brand System Design Guide 기반)"""
    RED = RGBColor(0xF3, 0x73, 0x21)            # = HANWHA_ORANGE (Primary accent, replaces LG RED)
    DARK_RED = RGBColor(0x35, 0x39, 0x68)       # = LIGHT_NAVY (Badge color, replaces LG DARK_RED)
    BLACK = RGBColor(0x1A, 0x1A, 0x2E)          # Text primary (navy-black)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)      # Text secondary
    MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)    # Text tertiary
    LIGHT_GRAY = RGBColor(0xEF, 0xEE, 0xE8)    # Surface (warm grey)
    BORDER_GRAY = RGBColor(0xC7, 0xBB, 0x9F)   # Sand border
    CHARCOAL = RGBColor(0x1D, 0x1E, 0x37)       # Navy (header bars)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN = RGBColor(0x2E, 0x7D, 0x32)
    ORANGE = RGBColor(0xF3, 0x73, 0x21)         # Hanwha Orange 100%
    BLUE = RGBColor(0x5C, 0x76, 0x87)           # Turquoise
    # Hanwha-specific extras
    NAVY = RGBColor(0x1D, 0x1E, 0x37)
    LIGHT_NAVY = RGBColor(0x35, 0x39, 0x68)
    ORANGE_70 = RGBColor(0xF8, 0x9B, 0x6C)
    ORANGE_50 = RGBColor(0xFB, 0xB5, 0x84)
    SAND = RGBColor(0xC7, 0xBB, 0x9F)
    TURQUOISE = RGBColor(0x5C, 0x76, 0x87)


THEMES = {
    "lg": {
        "colors": LGColors,
        "font_name": "에이투지체",
        "font_name_latin": "에이투지체",
        "fallback_font": "맑은 고딕",
        "cover_decoration": "l_bracket",
        "accent_bar_width": Cm(0.40),
        "separator_height": Cm(0.05),
    },
    "hanwha": {
        "colors": HanwhaColors,
        "font_name": "한화고딕",
        "font_name_latin": "한화고딕",
        "fallback_font": "맑은 고딕",
        "cover_decoration": "bottom_bar",
        "accent_bar_width": Cm(0.30),
        "separator_height": Cm(0.03),
    }
}


class LGTypography:
    """LG typography scale"""
    COVER_TITLE = Pt(32)
    COVER_SUBTITLE = Pt(14)
    SECTION_TITLE = Pt(28)       # Slide title (Bold) - tighter
    SUBTITLE = Pt(16)            # Slide subtitle (SemiBold) - tighter
    BODY_TITLE = Pt(16)
    BODY = Pt(12)                # All detail/body text (Regular) - denser
    BODY_SMALL = Pt(11)
    TABLE_HEADER = Pt(10)
    TABLE_BODY = Pt(9)
    CAPTION = Pt(9)
    TOC_TITLE = Pt(28)
    TOC_ITEM = Pt(16)
    TOC_SUBITEM = Pt(13)

    # ── Density presets ──
    # "spacious" — 인포그래픽/비전 슬라이드 (시각적 임팩트, 넉넉한 여백)
    BODY_SPACIOUS = Pt(16)
    BULLET_SPACIOUS = Pt(16)
    TABLE_HEADER_SPACIOUS = Pt(12)
    TABLE_BODY_SPACIOUS = Pt(11)
    CAPTION_SPACIOUS = Pt(10)

    # "compact" — 컨설팅 보고서 수준 (내용이 많은 슬라이드)
    BODY_COMPACT = Pt(10)
    BULLET_COMPACT = Pt(10)
    TABLE_HEADER_COMPACT = Pt(9)
    TABLE_BODY_COMPACT = Pt(8)
    CAPTION_COMPACT = Pt(8)

    # "dense" — 제안서/구축범위 수준 (빡빡한 테이블, 구성도)
    BODY_DENSE = Pt(9)
    BULLET_DENSE = Pt(9)
    TABLE_HEADER_DENSE = Pt(8)
    TABLE_BODY_DENSE = Pt(7)
    CAPTION_DENSE = Pt(7)


class LGDimensions:
    """Slide layout dimensions (16:9, 13.333 x 7.5 inches)"""
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Margins (report-dense layout)
    LEFT_MARGIN = Cm(0.8)
    RIGHT_MARGIN = Cm(0.5)
    TOP_MARGIN = Cm(0.5)
    BOTTOM_MARGIN = Cm(0.5)

    # NEW Chrome positions (from L-Style reference template)
    CHAPTER_LEFT = Cm(1.07)
    CHAPTER_TOP = Cm(0.78)
    CHAPTER_WIDTH = Cm(9.24)
    CHAPTER_HEIGHT = Cm(0.68)

    TITLE_LEFT = Cm(1.07)
    TITLE_TOP = Cm(2.16)
    TITLE_WIDTH = Cm(25.0)
    TITLE_HEIGHT = Cm(1.28)

    ACCENT_BAR_LEFT = Cm(0.67)
    ACCENT_BAR_TOP = Cm(2.30)
    ACCENT_BAR_WIDTH = Cm(0.40)
    ACCENT_BAR_HEIGHT = Cm(1.00)

    SUBTITLE_LEFT = Cm(1.07)
    SUBTITLE_TOP = Cm(3.74)
    SUBTITLE_WIDTH = Cm(20.0)
    SUBTITLE_HEIGHT = Cm(0.77)

    RED_LINE_LEFT = Cm(1.06)
    RED_LINE_TOP = Cm(4.87)
    RED_LINE_WIDTH = Cm(23.28)
    RED_LINE_HEIGHT = Cm(0.05)

    TAG_LEFT = Cm(23.89)
    TAG_TOP = Cm(0.30)
    TAG_WIDTH = Cm(9.00)
    TAG_HEIGHT = Cm(0.80)

    DOT_LEFT = Cm(33.07)
    DOT_TOP = Cm(0.50)
    DOT_SIZE = Cm(0.35)

    CONTENT_TOP = Cm(5.62)      # Content area starts here
    CONTENT_LEFT = Cm(1.33)
    CONTENT_RIGHT_MARGIN = Cm(0.5)

    # L-bracket (keep same)
    BRACKET_ARM_LENGTH = Cm(2.5)
    BRACKET_THICKNESS = Cm(0.4)


def _make_oxml_element(tag, **attribs):
    """Create an OxmlElement with attributes."""
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }
    element = etree.SubElement(etree.Element('dummy', nsmap=nsmap), qn(tag))
    for k, v in attribs.items():
        element.set(k, str(v))
    # Detach from dummy parent
    dummy = element.getparent()
    dummy.remove(element)
    return element


def _apply_shadow(shape, shadow_type="card"):
    """
    Apply template-matching shadow to a shape.
    
    Shadow types (from L-Style template analysis):
    - "card": WHITE container ambient shadow (blur=24.6pt, 11.2% opacity)
    - "content": F2F2F2 content box directional shadow (blur=4pt, dist=3pt, 45°, 40%)
    - "icon": Circle icon ambient shadow (blur=24.2pt, 14% opacity)
    """
    sp = shape._element
    spPr = sp.find(qn('a:spPr')) if hasattr(sp, 'find') else None
    if spPr is None:
        spPr = sp.spPr if hasattr(sp, 'spPr') else None
    if spPr is None:
        return
    
    # Remove existing effectLst
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    
    effectLst = _make_oxml_element('a:effectLst')
    
    if shadow_type == "card":
        # Ambient shadow: soft diffuse glow, no offset
        outerShdw = _make_oxml_element('a:outerShdw', blurRad='312382', algn='bl')
        srgbClr = _make_oxml_element('a:srgbClr', val='000000')
        alpha = _make_oxml_element('a:alpha', val='11221')
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
    elif shadow_type == "content":
        # Directional shadow: bottom-right, more visible
        outerShdw = _make_oxml_element('a:outerShdw',
            blurRad='50800', dist='38100', dir='2700000', algn='tl')
        srgbClr = _make_oxml_element('a:srgbClr', val='000000')
        alpha = _make_oxml_element('a:alpha', val='40000')
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
    elif shadow_type == "icon":
        # Icon ambient shadow: slightly more opaque than card
        outerShdw = _make_oxml_element('a:outerShdw', blurRad='307055', algn='bl')
        srgbClr = _make_oxml_element('a:srgbClr', val='000000')
        alpha = _make_oxml_element('a:alpha', val='14030')
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
    else:
        return
    
    effectLst.append(outerShdw)
    spPr.append(effectLst)



# ─────────────────────────────────────────────
# Main Builder Class
# ─────────────────────────────────────────────

class LGPresentation:
    """
    LG 그룹 브랜드 가이드라인을 따르는 PPTX 프레젠테이션 빌더 (멀티 테마 지원)

    Args:
        font_name: 한글 폰트 이름 (테마 기본값 사용 시 None)
        font_name_latin: 라틴 폰트 이름 (테마 기본값 사용 시 None)
        fallback_font: 폴백 폰트 (테마 기본값 사용 시 None)
        logo_path: 로고 이미지 파일 경로 (선택)
        theme: 테마 이름 ("lg" 또는 "hanwha", 기본: "lg")
    """

    def __init__(self, font_name=None, font_name_latin=None,
                 fallback_font=None, logo_path=None, theme="lg"):
        # Resolve theme configuration
        theme_config = THEMES.get(theme, THEMES["lg"])
        self.theme = theme
        self.colors = theme_config["colors"]

        self.font_name = font_name or theme_config["font_name"]
        self.font_name_latin = font_name_latin or theme_config["font_name_latin"]
        self.fallback_font = fallback_font or theme_config["fallback_font"]
        self.logo_path = logo_path

        # Theme-specific layout values
        self.accent_bar_width = theme_config["accent_bar_width"]
        self.separator_height = theme_config["separator_height"]

        # Weight-specific family names (per-theme font weight mapping)
        if theme == "hanwha":
            self.font_regular = "한화고딕 R_OTF"
            self.font_semibold = "한화 R_OTF"     # 중제목: 한화 오리지널
            self.font_bold = "한화 B_OTF"          # 대제목: 한화 오리지널 Bold
        else:
            # LG: each weight is a separate family (nameID=1)
            self.font_regular = f"{self.font_name} 4 Regular"
            self.font_semibold = f"{self.font_name} 6 SemiBold"
            self.font_bold = f"{self.font_name} 7 Bold"

        self.prs = Presentation()
        self.prs.slide_width = LGDimensions.SLIDE_WIDTH
        self.prs.slide_height = LGDimensions.SLIDE_HEIGHT

        # Patch the presentation theme so the default font resolves to
        # the correct weight-specific family name.
        self._patch_theme_fonts()

    # ─────────────────────────────────────────
    # Font Helpers
    # ─────────────────────────────────────────

    def _patch_theme_fonts(self):
        """
        Rewrite the presentation theme XML so that the default fonts resolve
        to the correct weight-specific family names.

        This font (에이투지체/A2Z) registers each weight as a separate family
        in its name table (e.g. "에이투지체 4 Regular", "에이투지체 7 Bold").
        PowerPoint on macOS uses these full family names (nameID=1), not the
        PrefFamily (nameID=16 = "에이투지체"). Without this theme patch,
        PowerPoint cannot resolve the font and falls back to Calibri/맑은 고딕.

        We set:
          - majorFont (headings) → bold weight
          - minorFont (body text) → regular weight
        """
        slide_master = self.prs.slide_masters[0]
        theme_part = None
        for rel in slide_master.part.rels.values():
            if 'theme' in rel.reltype:
                theme_part = rel.target_part
                break

        if theme_part is None:
            return

        theme_xml = etree.fromstring(theme_part.blob)
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

        font_map = {
            'a:majorFont': self.font_bold,     # headings → bold weight
            'a:minorFont': self.font_regular,   # body → regular weight
        }

        for font_group_tag, target_name in font_map.items():
            font_group = theme_xml.find(f'.//a:fontScheme/{font_group_tag}', ns)
            if font_group is None:
                continue

            for child_tag in ['a:latin', 'a:ea', 'a:cs']:
                child = font_group.find(child_tag, ns)
                if child is not None:
                    child.set('typeface', target_name)

            for font_el in font_group.findall('a:font', ns):
                if font_el.get('script') == 'Hang':
                    font_el.set('typeface', target_name)

        theme_part._blob = etree.tostring(
            theme_xml, xml_declaration=True, encoding='UTF-8', standalone=True
        )

    def _set_font(self, run, size=None, bold=False, italic=False, color=None,
                  weight=None):
        """
        Set font properties including East Asian font for Korean text.

        Args:
            weight: Explicit weight override — "bold", "semibold", or "regular".
                    If omitted, falls back to ``bold`` flag (True→bold, False→regular).
        """
        # Pick the weight-specific family name.
        # This font registers each weight as a separate family (nameID=1):
        #   "에이투지체 7 Bold", "에이투지체 6 SemiBold", "에이투지체 4 Regular"
        # PowerPoint on macOS resolves via these full names, so we must
        # use the correct one. We never set font.bold=True because the
        # weight is already embedded in the family name.
        if weight == "semibold":
            target_font = self.font_semibold
        elif weight == "bold" or (weight is None and bold):
            target_font = self.font_bold
        else:
            target_font = self.font_regular

        font = run.font
        if size:
            font.size = size
        font.bold = False          # never synthetic bold; weight is in the family name
        font.italic = italic
        if color:
            font.color.rgb = color

        # Directly manipulate XML for reliable font binding
        rPr = run._r.get_or_add_rPr()

        # Set language attributes on the run properties element
        rPr.set('lang', 'ko-KR')
        rPr.set('altLang', 'en-US')

        # Remove existing font elements to avoid duplicates
        for tag in ['a:latin', 'a:ea', 'a:cs']:
            for existing in rPr.findall(qn(tag)):
                rPr.remove(existing)

        # <a:latin> — Latin font
        latin = _make_oxml_element('a:latin')
        latin.set('typeface', target_font)
        rPr.append(latin)

        # <a:ea> — East Asian font (charset="-127" for CJK binding)
        ea = _make_oxml_element('a:ea')
        ea.set('typeface', target_font)
        ea.set('charset', '-127')
        rPr.append(ea)

        # <a:cs> — Complex Script font
        cs = _make_oxml_element('a:cs')
        cs.set('typeface', target_font)
        cs.set('charset', '-127')
        rPr.append(cs)

    def _add_text(self, text_frame, text, size=None, bold=False, color=None,
                  alignment=None, space_before=None, space_after=None,
                  weight=None):
        """Add a paragraph with formatted text to a text frame."""
        if text_frame.paragraphs and text_frame.paragraphs[0].text == '':
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        if alignment:
            p.alignment = alignment
        if space_before:
            p.space_before = space_before
        if space_after:
            p.space_after = space_after

        run = p.add_run()
        run.text = text
        self._set_font(run, size=size, bold=bold, color=color, weight=weight)
        return p

    def _add_textbox(self, slide, left, top, width, height, text="",
                     size=None, bold=False, color=None, alignment=None,
                     word_wrap=True, vertical=None, weight=None):
        """Add a text box with formatted text."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        if vertical:
            tf.vertical_anchor = vertical

        if text:
            self._add_text(tf, text, size=size, bold=bold, color=color,
                          alignment=alignment, weight=weight)

        # Make transparent
        txBox.fill.background()
        txBox.line.fill.background()
        return txBox

    # ─────────────────────────────────────────
    # Visual Element Helpers
    # ─────────────────────────────────────────

    def _add_l_bracket(self, slide, corner="top-left", arm_length=None,
                       thickness=None, color=None):
        """
        Draw an L-bracket corner decoration.

        Args:
            corner: "top-left" or "bottom-right"
            arm_length: Length of bracket arms (default: 2.5cm)
            thickness: Thickness of bracket arms (default: 0.4cm)
            color: Fill color (default: LG RED)
        """
        arm = arm_length or LGDimensions.BRACKET_ARM_LENGTH
        thick = thickness or LGDimensions.BRACKET_THICKNESS
        clr = color or self.colors.RED

        shapes = slide.shapes
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)
        a = int(arm)
        t = int(thick)

        if corner == "top-left":
            # Horizontal arm (top)
            h_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Cm(0.8), Cm(0.8), a, t)
            h_bar.fill.solid()
            h_bar.fill.fore_color.rgb = clr
            h_bar.line.fill.background()

            # Vertical arm (left)
            v_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Cm(0.8), Cm(0.8), t, a)
            v_bar.fill.solid()
            v_bar.fill.fore_color.rgb = clr
            v_bar.line.fill.background()

        elif corner == "bottom-right":
            # Position from bottom-right
            br_x = sw - int(Cm(0.8)) - a
            br_y = sh - int(Cm(0.8)) - a

            # Horizontal arm (bottom)
            h_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     br_x, br_y + a - t, a, t)
            h_bar.fill.solid()
            h_bar.fill.fore_color.rgb = clr
            h_bar.line.fill.background()

            # Vertical arm (right)
            v_bar = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     br_x + a - t, br_y, t, a)
            v_bar.fill.solid()
            v_bar.fill.fore_color.rgb = clr
            v_bar.line.fill.background()

    def _add_content_chrome(self, slide, title="", section="", subtitle="",
                            chapter=""):
        """
        Add all L-Style chrome elements for a content slide.
        Returns the y position where content should start (below red line).

        Chrome layout (from reference template):
          - Chapter name: top-left, 10pt RED bold
          - Accent bar: left of title, RED solid
          - Title: 24pt BLACK bold
          - Subtitle: 16pt #333 SemiBold
          - Red separator line: full-width RED
          - Tag text: top-right, 11pt #666 right-aligned
          - Red dot: far top-right, RED oval
        """
        # 1. Chapter name (top-left, small red text)
        if chapter:
            self._add_textbox(
                slide,
                LGDimensions.CHAPTER_LEFT, LGDimensions.CHAPTER_TOP,
                LGDimensions.CHAPTER_WIDTH, LGDimensions.CHAPTER_HEIGHT,
                text=chapter, size=Pt(10), bold=True, color=self.colors.RED
            )

        # 2. Accent bar (left of title) — width from theme config
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            LGDimensions.ACCENT_BAR_LEFT, LGDimensions.ACCENT_BAR_TOP,
            self.accent_bar_width, LGDimensions.ACCENT_BAR_HEIGHT
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors.RED
        bar.line.fill.background()

        # 3. Title
        if title:
            self._add_textbox(
                slide,
                LGDimensions.TITLE_LEFT, LGDimensions.TITLE_TOP,
                LGDimensions.TITLE_WIDTH, LGDimensions.TITLE_HEIGHT,
                text=title, size=Pt(24), bold=True, color=self.colors.BLACK,
                vertical=MSO_ANCHOR.MIDDLE
            )

        # 4. Subtitle
        if subtitle:
            self._add_textbox(
                slide,
                LGDimensions.SUBTITLE_LEFT, LGDimensions.SUBTITLE_TOP,
                LGDimensions.SUBTITLE_WIDTH, LGDimensions.SUBTITLE_HEIGHT,
                text=subtitle, size=Pt(16), color=self.colors.DARK_GRAY,
                weight="semibold"
            )

        # 5. Red separator line — height from theme config
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            LGDimensions.RED_LINE_LEFT, LGDimensions.RED_LINE_TOP,
            LGDimensions.RED_LINE_WIDTH, self.separator_height
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors.RED
        line.line.fill.background()

        # 6. Tag text (top-right)
        if section:
            self._add_textbox(
                slide,
                LGDimensions.TAG_LEFT, LGDimensions.TAG_TOP,
                LGDimensions.TAG_WIDTH, LGDimensions.TAG_HEIGHT,
                text=section, size=Pt(11), color=self.colors.MEDIUM_GRAY,
                alignment=PP_ALIGN.RIGHT
            )

        # 7. Red dot (far top-right)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            LGDimensions.DOT_LEFT, LGDimensions.DOT_TOP,
            LGDimensions.DOT_SIZE, LGDimensions.DOT_SIZE
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = self.colors.RED
        dot.line.fill.background()

        return LGDimensions.CONTENT_TOP  # content start y

    def _add_horizontal_line(self, slide, left, top, width, color=None, height=None):
        """Add a thin horizontal line."""
        clr = color or self.colors.BORDER_GRAY
        h = height or Pt(1)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, h
        )
        line.fill.solid()
        line.fill.fore_color.rgb = clr
        line.line.fill.background()
        return line

    def _add_shadow(self, shape, shadow_type="card"):
        """
        Add a shadow to a shape matching the L-Style template exactly.

        Shadow types (from template analysis):
        - "card": Ambient glow shadow for WHITE card containers
                  (blur=312382 EMU, alpha=11.2%, no offset)
        - "icon": Ambient shadow for circle icons
                  (blur=307055 EMU, alpha=14%, no offset)
        - "content": Directional shadow for secondary content boxes
                     (blur=50800, dist=38100, dir=2700000, alpha=40%)
        """
        spPr = shape._element.spPr

        # Remove any existing effectLst
        for existing in spPr.findall(qn('a:effectLst')):
            spPr.remove(existing)

        effectLst = _make_oxml_element('a:effectLst')

        if shadow_type == "card":
            outerShdw = _make_oxml_element('a:outerShdw')
            outerShdw.set('blurRad', '312382')
            outerShdw.set('algn', 'bl')
            srgbClr = _make_oxml_element('a:srgbClr')
            srgbClr.set('val', '000000')
            alphaElem = _make_oxml_element('a:alpha')
            alphaElem.set('val', '11221')
            srgbClr.append(alphaElem)
            outerShdw.append(srgbClr)
        elif shadow_type == "icon":
            outerShdw = _make_oxml_element('a:outerShdw')
            outerShdw.set('blurRad', '307055')
            outerShdw.set('algn', 'bl')
            srgbClr = _make_oxml_element('a:srgbClr')
            srgbClr.set('val', '000000')
            alphaElem = _make_oxml_element('a:alpha')
            alphaElem.set('val', '14030')
            srgbClr.append(alphaElem)
            outerShdw.append(srgbClr)
        elif shadow_type == "content":
            outerShdw = _make_oxml_element('a:outerShdw')
            outerShdw.set('blurRad', '50800')
            outerShdw.set('dist', '38100')
            outerShdw.set('dir', '2700000')
            outerShdw.set('algn', 'tl')
            outerShdw.set('rotWithShape', '0')
            srgbClr = _make_oxml_element('a:srgbClr')
            srgbClr.set('val', '000000')
            alphaElem = _make_oxml_element('a:alpha')
            alphaElem.set('val', '40000')
            srgbClr.append(alphaElem)
            outerShdw.append(srgbClr)
        else:
            return

        effectLst.append(outerShdw)
        spPr.append(effectLst)

    def _get_blank_slide(self):
        """Get a blank slide layout and add a new slide."""
        # Try to find blank layout (usually index 6, but may vary)
        layout = self.prs.slide_layouts[6]  # Blank
        return self.prs.slides.add_slide(layout)

    # ─────────────────────────────────────────
    # Table Helpers
    # ─────────────────────────────────────────

    def _set_cell_border(self, cell, color="CCCCCC", width='6350'):
        """Set all four borders on a table cell."""
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        for edge in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            for existing in tcPr.findall(qn(edge)):
                tcPr.remove(existing)
            ln = _make_oxml_element(edge)
            ln.set('w', str(width))
            ln.set('cap', 'flat')
            ln.set('cmpd', 'sng')
            ln.set('algn', 'ctr')

            solidFill = _make_oxml_element('a:solidFill')
            srgbClr = _make_oxml_element('a:srgbClr')
            srgbClr.set('val', color)
            solidFill.append(srgbClr)
            ln.append(solidFill)

            prstDash = _make_oxml_element('a:prstDash')
            prstDash.set('val', 'solid')
            ln.append(prstDash)

            tcPr.append(ln)

    def _format_table_cell(self, cell, text="", is_header=False, font_size=None,
                           alignment=None, bold=None, text_color=None, bg_color=None,
                           density="normal"):
        """Format a table cell with LG styling.

        Args:
            density: "normal" | "compact" | "dense"
                     — controls default font size and cell margins.
        """
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = alignment or PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text

        # Resolve default font size per density
        if font_size is None:
            if density == "spacious":
                default_hdr = LGTypography.TABLE_HEADER_SPACIOUS
                default_body = LGTypography.TABLE_BODY_SPACIOUS
            elif density == "compact":
                default_hdr = LGTypography.TABLE_HEADER_COMPACT
                default_body = LGTypography.TABLE_BODY_COMPACT
            elif density == "dense":
                default_hdr = LGTypography.TABLE_HEADER_DENSE
                default_body = LGTypography.TABLE_BODY_DENSE
            else:
                default_hdr = LGTypography.TABLE_HEADER
                default_body = LGTypography.TABLE_BODY
        else:
            default_hdr = font_size
            default_body = font_size

        if is_header:
            self._set_font(run, size=default_hdr,
                          bold=True if bold is None else bold,
                          color=text_color or self.colors.WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color or self.colors.CHARCOAL
        else:
            self._set_font(run, size=default_body,
                          bold=bold or False,
                          color=text_color or self.colors.BLACK)
            if bg_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
            else:
                cell.fill.background()

        # Margins — tighter for compact/dense, wider for spacious
        if density == "spacious":
            cell.margin_left = Cm(0.3)
            cell.margin_right = Cm(0.3)
            cell.margin_top = Cm(0.15)
            cell.margin_bottom = Cm(0.15)
        elif density == "dense":
            cell.margin_left = Cm(0.1)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.05)
            cell.margin_bottom = Cm(0.05)
        elif density == "compact":
            cell.margin_left = Cm(0.15)
            cell.margin_right = Cm(0.15)
            cell.margin_top = Cm(0.08)
            cell.margin_bottom = Cm(0.08)
        else:
            cell.margin_left = Cm(0.2)
            cell.margin_right = Cm(0.2)
            cell.margin_top = Cm(0.1)
            cell.margin_bottom = Cm(0.1)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ─────────────────────────────────────────
    # Slide Methods
    # ─────────────────────────────────────────

    def add_cover(self, title, subtitle="", date="", logo_path=None):
        """
        Add a cover slide with theme-specific decorations.

        Args:
            title: Main presentation title
            subtitle: Team/department name
            date: Presentation date (e.g., "2025.10.23")
            logo_path: Override logo path for this slide
        """
        slide = self._get_blank_slide()

        # Theme-specific cover decoration
        if self.theme == "hanwha":
            # Hanwha: bottom orange accent bar instead of L-brackets
            bar_h = Cm(0.15)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Cm(0), LGDimensions.SLIDE_HEIGHT - bar_h,
                LGDimensions.SLIDE_WIDTH, bar_h
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.colors.RED
            bar.line.fill.background()
        else:
            # LG: L-brackets
            self._add_l_bracket(slide, "top-left")
            self._add_l_bracket(slide, "bottom-right")

        # Title (centered)
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        title_box = self._add_textbox(
            slide,
            left=Cm(3), top=sh // 2 - Cm(3),
            width=sw - Cm(6), height=Cm(4),
            text=title,
            size=LGTypography.COVER_TITLE,
            bold=True,
            color=self.colors.BLACK,
            alignment=PP_ALIGN.CENTER,
            vertical=MSO_ANCHOR.MIDDLE
        )

        # Subtitle + Date (bottom center)
        info_parts = []
        if subtitle:
            info_parts.append(subtitle)
        if date:
            info_parts.append(date)

        if info_parts:
            info_box = self._add_textbox(
                slide,
                left=Cm(3), top=sh - Cm(4),
                width=sw - Cm(6), height=Cm(2.5),
                alignment=PP_ALIGN.CENTER
            )
            tf = info_box.text_frame
            tf.clear()

            if subtitle:
                p1 = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                p1.alignment = PP_ALIGN.CENTER
                run1 = p1.add_run()
                run1.text = subtitle
                self._set_font(run1, size=LGTypography.COVER_SUBTITLE,
                              color=self.colors.DARK_GRAY)

            if date:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(4)
                run2 = p2.add_run()
                run2.text = date
                self._set_font(run2, size=LGTypography.COVER_SUBTITLE,
                              bold=True, color=self.colors.BLACK)

        # Logo (if provided)
        logo = logo_path or self.logo_path
        if logo and os.path.exists(logo):
            slide.shapes.add_picture(
                logo,
                sw - Cm(5), sh - Cm(4),
                width=Cm(3)
            )

        return slide

    def add_toc(self, items):
        """
        Add a Table of Contents slide.

        Args:
            items: List of tuples (title, [sub_items])
                   e.g., [("Summary", []), ("시스템 소개", ["항목1", "항목2"])]
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # Top gray line
        self._add_horizontal_line(
            slide, Cm(1.5), Cm(1.5), sw - Cm(3)
        )

        # "Contents" title
        contents_box = self._add_textbox(
            slide,
            left=Cm(2), top=Cm(2.2),
            width=Cm(10), height=Cm(1.5),
            text="Contents",
            size=LGTypography.TOC_TITLE,
            color=self.colors.BLACK
        )

        # Red underline bar
        red_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(2), Cm(4.0), Cm(2.5), Cm(0.15)
        )
        red_bar.fill.solid()
        red_bar.fill.fore_color.rgb = self.colors.RED
        red_bar.line.fill.background()

        # Gray divider line below title
        self._add_horizontal_line(
            slide, Cm(2), Cm(4.5), sw - Cm(4)
        )

        # TOC items
        y_pos = Cm(5.0)
        roman = ['I', 'II', 'III', 'IV', 'V', 'VI', 'VII', 'VIII', 'IX', 'X']

        for i, (title, sub_items) in enumerate(items):
            # Roman numeral + title
            numeral = roman[i] if i < len(roman) else str(i + 1)
            toc_text = f"{numeral}.   {title}"

            toc_box = self._add_textbox(
                slide,
                left=Cm(3), top=y_pos,
                width=Cm(20), height=Cm(0.9),
                text=toc_text,
                size=LGTypography.TOC_ITEM,
                bold=True,
                color=self.colors.RED
            )
            y_pos += Cm(0.85)

            # Sub-items
            if sub_items:
                for sub in sub_items:
                    sub_box = self._add_textbox(
                        slide,
                        left=Cm(4.5), top=y_pos,
                        width=Cm(20), height=Cm(0.6),
                        text=f"- {sub}",
                        size=LGTypography.TOC_SUBITEM,
                        color=self.colors.DARK_GRAY
                    )
                    y_pos += Cm(0.55)

            y_pos += Cm(0.2)

        return slide

    def add_section_divider(self, number, title, color=None):
        """
        Add a modern section divider slide.

        Clean, minimal design: large section number with accent line,
        title text, and subtle decorative elements.

        Args:
            number: Section number (e.g., "I", "II", "01")
            title: Section title
            color: Override accent color (default: LG RED)
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)
        accent = color or self.colors.RED

        # ── Background: clean white ──
        # Subtle bottom accent bar (thin, full width)
        bottom_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, sh - int(Cm(0.4)), sw, int(Cm(0.4))
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = accent
        bottom_bar.line.fill.background()

        # ── Left vertical accent line ──
        vert_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(Cm(3.5)), int(Cm(5.0)),
            int(Cm(0.25)), int(Cm(8.5))
        )
        vert_line.fill.solid()
        vert_line.fill.fore_color.rgb = accent
        vert_line.line.fill.background()

        # ── Section number (large, light gray) ──
        num_box = self._add_textbox(
            slide,
            left=Cm(5.0), top=Cm(4.0),
            width=Cm(10), height=Cm(5.5),
            text=str(number),
            size=Pt(72),
            bold=True,
            color=RGBColor(0xE0, 0xE0, 0xE0),
            alignment=PP_ALIGN.LEFT
        )

        # ── Title (bold, overlapping number area) ──
        self._add_textbox(
            slide,
            left=Cm(5.0), top=Cm(9.0),
            width=Cm(24), height=Cm(3.5),
            text=title,
            size=Pt(32),
            bold=True,
            color=self.colors.BLACK,
            alignment=PP_ALIGN.LEFT,
            vertical=MSO_ANCHOR.TOP
        )

        # ── Short red underline below title ──
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(Cm(5.0)), int(Cm(12.8)),
            int(Cm(4.0)), int(Cm(0.15))
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = accent
        underline.line.fill.background()

        return slide

    def add_content(self, title, section="", body="", bullets=None,
                    sub_title="", subtitle="", chapter="",
                    density="normal"):
        """
        Add a standard content slide with L-Style chrome.

        Args:
            title: Slide title (e.g., "1.1 시스템 개요")
            section: Section name shown in top-right tag
            body: Body text paragraph
            bullets: List of bullet point strings
            sub_title: Subtitle text below title (legacy alias)
            subtitle: Subtitle text below title (preferred)
            chapter: Chapter name shown top-left in red
            density: "normal" | "compact" | "dense"
                     — controls font size and spacing for body/bullets.
        """
        slide = self._get_blank_slide()
        _subtitle = subtitle or sub_title

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=_subtitle, chapter=chapter)
        y_pos = int(y_pos)

        # Resolve sizes per density
        if density == "spacious":
            body_size = LGTypography.BODY_SPACIOUS
            bullet_size_few = LGTypography.BODY_SPACIOUS
            bullet_size_many = LGTypography.BULLET_SPACIOUS
            space_few = Pt(14)
            space_many = Pt(8)
        elif density == "dense":
            body_size = LGTypography.BODY_DENSE
            bullet_size_few = LGTypography.BODY_DENSE
            bullet_size_many = LGTypography.BULLET_DENSE
            space_few = Pt(4)
            space_many = Pt(2)
        elif density == "compact":
            body_size = LGTypography.BODY_COMPACT
            bullet_size_few = LGTypography.BODY_COMPACT
            bullet_size_many = LGTypography.BULLET_COMPACT
            space_few = Pt(6)
            space_many = Pt(2)
        else:
            body_size = Pt(14)
            bullet_size_few = Pt(14)
            bullet_size_many = LGTypography.BODY
            space_few = Pt(8)
            space_many = Pt(3)

        # Body text
        if body:
            body_box = self._add_textbox(
                slide,
                left=LGDimensions.CONTENT_LEFT, top=y_pos,
                width=Cm(31), height=Cm(12),
                text=body,
                size=body_size,
                color=self.colors.BLACK
            )
            y_pos += int(Cm(2.0))

        # Bullet points
        if bullets:
            few_bullets = len(bullets) <= 4 and not body

            bullet_box = slide.shapes.add_textbox(
                LGDimensions.CONTENT_LEFT, y_pos, Cm(31), Cm(12)
            )
            tf = bullet_box.text_frame
            tf.word_wrap = True
            bullet_box.fill.background()
            bullet_box.line.fill.background()

            bullet_font_size = bullet_size_few if few_bullets else bullet_size_many
            bullet_space_after = space_few if few_bullets else space_many

            for j, bullet in enumerate(bullets):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.space_after = bullet_space_after

                # Bullet character
                pPr = p._pPr
                if pPr is None:
                    pPr = p._p.get_or_add_pPr()
                buChar = _make_oxml_element('a:buChar')
                buChar.set('char', '\u2022')
                pPr.append(buChar)

                # Indent
                pPr.set('marL', str(int(Cm(0.5))))
                pPr.set('indent', str(int(Cm(-0.4))))

                run = p.add_run()
                run.text = bullet
                self._set_font(run, size=bullet_font_size, color=self.colors.BLACK)

        return slide

    def add_roadmap(self, title, section="", subtitle="", years=None,
                    roadmap_items=None, table_data=None, chapter=""):
        """
        Add a roadmap timeline slide.

        Args:
            title: Slide title (e.g., "[프로젝트] 로드맵 및 계열사별 비교")
            section: Section name for top-right indicator
            subtitle: Description text below title
            years: List of year/phase strings for timeline headers
                   e.g., ["(2025) 데이터 안정화", "(2026) 분석 고도화", "(2027) 자동화"]
            roadmap_items: Dict with structure:
                {
                    "label": "시스템 로드맵",  # Left label text
                    "rows": [
                        {
                            "items_by_year": [
                                # Items for each year column
                                [{"text": "항목", "tag": "계열사A", "tag_color": "#1565C0"}],
                                [{"text": "항목2", "color": "green"}],
                                [{"text": "항목3", "color": "orange"}]
                            ]
                        }
                    ]
                }
            table_data: Dict with structure:
                {
                    "title": "계열사 별 현황",
                    "headers": ["계열사", "DX 수준", "AI 적용", ...],
                    "rows": [["에너지솔루션", "Lv3", "적용(예정)", ...], ...]
                }
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        # Timeline headers (chevron/arrow style)
        if years:
            num_years = len(years)
            header_left = Cm(4.5)
            total_width = sw - int(header_left) - int(Cm(1.5))
            col_width = total_width // num_years

            for i, year_text in enumerate(years):
                x = int(header_left) + (col_width * i)
                header_bar = slide.shapes.add_shape(
                    MSO_SHAPE.CHEVRON,
                    x, y_pos, col_width, Cm(1.2)
                )
                header_bar.fill.solid()
                header_bar.fill.fore_color.rgb = self.colors.CHARCOAL
                header_bar.line.fill.background()
                self._add_shadow(header_bar)

                # Year text inside header
                tf = header_bar.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = year_text
                self._set_font(run, size=Pt(11), bold=True, color=self.colors.WHITE)

            y_pos += int(Cm(1.8))

        # Roadmap content grid
        if roadmap_items:
            label_text = roadmap_items.get("label", "")
            rows = roadmap_items.get("rows", [])

            # Left label block
            if label_text and years:
                label_height = Cm(len(rows) * 4.5) if rows else Cm(4)
                label_block = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Cm(1.5), y_pos, Cm(2.5), label_height
                )
                label_block.fill.solid()
                label_block.fill.fore_color.rgb = self.colors.RED
                label_block.line.fill.background()
                self._add_shadow(label_block)

                # Label text (vertical-ish)
                tf = label_block.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                for char_idx, line in enumerate(label_text.split('\n') if '\n' in label_text else [label_text]):
                    if char_idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = line
                    self._set_font(run, size=Pt(11), bold=True, color=self.colors.WHITE)

            # Content cells
            if years and rows:
                content_left = Cm(4.5)
                total_width = sw - int(content_left) - int(Cm(1.5))
                col_width = total_width // len(years)

                for row_idx, row_data in enumerate(rows):
                    items_by_year = row_data.get("items_by_year", [])
                    row_top = y_pos + (int(Cm(4.5)) * row_idx)

                    for col_idx, year_items in enumerate(items_by_year):
                        cell_x = int(content_left) + (col_width * col_idx)

                        # Cell background (WHITE rounded rectangle)
                        cell_bg = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            cell_x + Cm(0.1), row_top + Cm(0.1),
                            col_width - Cm(0.2), Cm(4.0)
                        )
                        cell_bg.fill.solid()
                        cell_bg.fill.fore_color.rgb = self.colors.WHITE
                        cell_bg.line.color.rgb = self.colors.RED
                        cell_bg.line.width = Pt(1)
                        # Set corner radius (template: val 8000)
                        _spPr = cell_bg._element.spPr
                        _pg = _spPr.find(qn('a:prstGeom'))
                        if _pg is not None:
                            _al = _pg.find(qn('a:avLst'))
                            if _al is None:
                                _al = _make_oxml_element('a:avLst')
                                _pg.append(_al)
                            for _gd in _al.findall(qn('a:gd')):
                                _al.remove(_gd)
                            _gd = _make_oxml_element('a:gd')
                            _gd.set('name', 'adj')
                            _gd.set('fmla', 'val 8000')
                            _al.append(_gd)
                        self._add_shadow(cell_bg)

                        # Cell content
                        if year_items:
                            cell_box = slide.shapes.add_textbox(
                                cell_x + Cm(0.4), row_top + Cm(0.4),
                                col_width - Cm(0.8), Cm(3.5)
                            )
                            cell_tf = cell_box.text_frame
                            cell_tf.word_wrap = True
                            cell_box.fill.background()
                            cell_box.line.fill.background()

                            for item_idx, item in enumerate(year_items):
                                if item_idx == 0:
                                    p = cell_tf.paragraphs[0]
                                else:
                                    p = cell_tf.add_paragraph()

                                p.space_after = Pt(3)

                                # Determine text color
                                item_text = item if isinstance(item, str) else item.get("text", "")
                                item_color_name = None if isinstance(item, str) else item.get("color", None)

                                if item_color_name == "green":
                                    text_color = self.colors.GREEN
                                elif item_color_name == "orange":
                                    text_color = self.colors.ORANGE
                                else:
                                    text_color = self.colors.BLACK

                                # Bullet
                                pPr = p._p.get_or_add_pPr()
                                buChar = _make_oxml_element('a:buChar')
                                buChar.set('char', '\u2022')
                                pPr.append(buChar)

                                run = p.add_run()
                                run.text = item_text
                                self._set_font(run, size=LGTypography.BODY_SMALL,
                                              color=text_color)

                                # Tag badge (e.g., "계열사A")
                                if isinstance(item, dict) and item.get("tag"):
                                    tag_run = p.add_run()
                                    tag_run.text = f"  {item['tag']}"
                                    tag_color_hex = item.get("tag_color", "#1565C0")
                                    r, g, b = int(tag_color_hex[1:3], 16), int(tag_color_hex[3:5], 16), int(tag_color_hex[5:7], 16)
                                    self._set_font(tag_run, size=Pt(9), bold=True,
                                                  color=RGBColor(r, g, b))

        # Table section
        if table_data:
            self._add_roadmap_table(slide, table_data, y_start=y_pos)

        return slide

    def _add_roadmap_table(self, slide, table_data, y_start=None):
        """Add a comparison table at the bottom of a slide."""
        title = table_data.get("title", "")
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        if not headers or not rows:
            return

        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # Table title
        table_y = sh - Cm(1.5) - Cm(len(rows) * 1.0 + 1.5)
        if title:
            self._add_textbox(
                slide,
                left=Cm(1.5), top=table_y - Cm(1.2),
                width=Cm(15), height=Cm(1.0),
                text=f"\u25a0 {title}",
                size=LGTypography.BODY,
                bold=True,
                color=self.colors.BLACK
            )

        # Table
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        table_width = sw - int(Cm(3))
        table_height = Cm(num_rows * 1.0)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Cm(1.5), table_y,
            table_width, table_height
        )
        table = table_shape.table

        # Disable default styling
        table.first_row = False
        table.first_col = False
        table.last_row = False
        table.last_col = False
        table.horz_banding = False
        table.vert_banding = False

        # Set column widths evenly
        col_width = table_width // num_cols
        for c in range(num_cols):
            table.columns[c].width = col_width

        # Header row
        for c, header in enumerate(headers):
            self._format_table_cell(table.cell(0, c), header, is_header=True)
            self._set_cell_border(table.cell(0, c), color="3C3C3C")

        # Data rows
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                bg = self.colors.WHITE if r % 2 == 0 else self.colors.LIGHT_GRAY
                self._format_table_cell(table.cell(r + 1, c), str(val), bg_color=bg)
                self._set_cell_border(table.cell(r + 1, c))

    def add_table(self, title, section="", headers=None, rows=None,
                  col_widths=None, subtitle="", chapter="",
                  density="normal", merge_column=None,
                  row_alignment=None, header_alignment=None):
        """
        Add a dedicated table slide.

        Args:
            title: Slide title
            section: Section name for top-right
            headers: List of column header strings
            rows: List of row data (list of lists).
                  To merge cells vertically, set the merge column value
                  to "" for continuation rows (same as the row above).
            col_widths: Optional list of column widths in Cm
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
            density: "normal" | "compact" | "dense"
                     — controls font size, cell margins, row height
            merge_column: Column index (0-based) to auto-merge vertically
                          when consecutive cells have the same value.
                          Can be int or list of ints for multiple columns.
            row_alignment: PP_ALIGN value for data cells (default: CENTER)
            header_alignment: PP_ALIGN value for header cells (default: CENTER)
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        if not headers or not rows:
            return slide

        # Table dimensions — adjust per density level
        num_rows = len(rows) + 1
        num_cols = len(headers)

        if density == "spacious":
            table_left = int(Cm(2.0))
            table_width = sw - int(Cm(3.0))
            bottom_margin = int(Cm(0.8))
            max_row_height = Cm(1.6)
        elif density == "dense":
            table_left = int(Cm(1.0))
            table_width = sw - int(Cm(1.5))
            bottom_margin = int(Cm(0.3))
            max_row_height = Cm(0.85)
        elif density == "compact":
            table_left = int(Cm(1.2))
            table_width = sw - int(Cm(1.7))
            bottom_margin = int(Cm(0.4))
            max_row_height = Cm(1.0)
        else:
            table_left = int(Cm(1.5))
            table_width = sw - int(Cm(2.0))
            bottom_margin = int(Cm(0.5))
            max_row_height = Cm(1.2)

        available_height = int(LGDimensions.SLIDE_HEIGHT) - y_pos - bottom_margin
        row_height = min(max_row_height, available_height // num_rows)
        table_height = row_height * num_rows

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            table_left, y_pos,
            table_width, table_height
        )
        table = table_shape.table

        # Disable default styling
        table.first_row = False
        table.first_col = False
        table.last_row = False
        table.last_col = False
        table.horz_banding = False
        table.vert_banding = False

        # Set column widths
        if col_widths:
            for c, w in enumerate(col_widths):
                if c < num_cols:
                    table.columns[c].width = Cm(w)
        else:
            col_width = table_width // num_cols
            for c in range(num_cols):
                table.columns[c].width = col_width

        # Header row
        for c, header in enumerate(headers):
            h_align = header_alignment or PP_ALIGN.CENTER
            self._format_table_cell(table.cell(0, c), header, is_header=True,
                                    density=density, alignment=h_align)
            self._set_cell_border(table.cell(0, c), color="3C3C3C")

        # Data rows
        d_align = row_alignment or PP_ALIGN.CENTER
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                bg = self.colors.WHITE if r % 2 == 0 else self.colors.LIGHT_GRAY
                self._format_table_cell(table.cell(r + 1, c), str(val),
                                        bg_color=bg, density=density,
                                        alignment=d_align)
                self._set_cell_border(table.cell(r + 1, c))

        # Vertical merge for specified columns
        if merge_column is not None:
            merge_cols = merge_column if isinstance(merge_column, (list, tuple)) else [merge_column]
            for mc in merge_cols:
                if mc >= num_cols:
                    continue
                r = 1  # start from first data row (row index 1 in table)
                while r < num_rows:
                    cell_text = table.cell(r, mc).text.strip()
                    if not cell_text:
                        r += 1
                        continue
                    # Find consecutive rows with same value or empty string
                    merge_end = r + 1
                    while merge_end < num_rows:
                        next_text = table.cell(merge_end, mc).text.strip()
                        if next_text == "" or next_text == cell_text:
                            merge_end += 1
                        else:
                            break
                    if merge_end > r + 1:
                        # Merge cells vertically
                        top_cell = table.cell(r, mc)
                        bottom_cell = table.cell(merge_end - 1, mc)
                        top_cell.merge(bottom_cell)
                        # Re-format the merged cell
                        bg = self.colors.LIGHT_GRAY
                        self._format_table_cell(top_cell, cell_text,
                                                bg_color=bg, density=density,
                                                bold=True, alignment=PP_ALIGN.CENTER)
                        self._set_cell_border(top_cell)
                        r = merge_end
                    else:
                        r += 1

        return slide

    def add_blank_content(self, title, section="", subtitle="", chapter=""):
        """
        Add a blank content slide with just the L-Style chrome.
        Returns (slide, y_pos) tuple for custom free-form content addition.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red

        Returns:
            tuple: (slide, y_pos) where y_pos is the pixel position
                   where content should start (below red separator line).
        """
        slide = self._get_blank_slide()

        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)

        return slide, int(y_pos)

    def add_recommendation(self, title="제언", section="", subtitle="",
                           recommendations=None, chapter=""):
        """
        Add a recommendation slide — commonly placed before
        the closing slide. Displays numbered recommendation items
        with optional detail text.

        Args:
            title: Slide title (default: "제언")
            section: Section name for top-right
            subtitle: Optional subtitle
            recommendations: List of strings or list of dicts:
                ["데이터 품질 확보가 최우선", "단계적 MLOps 성숙도 향상"]
                or
                [
                    {"title": "데이터 품질 확보", "detail": "데이터 정합성 검증 체계 구축"},
                    {"title": "MLOps 성숙도 향상", "detail": "단계적 자동화 파이프라인 도입"},
                ]
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        recommendations = recommendations or []
        if not recommendations:
            return slide

        # Numbered recommendation items
        item_y = y_pos + int(Cm(0.3))
        for idx, rec in enumerate(recommendations):
            if isinstance(rec, dict):
                rec_title = rec.get("title", "")
                rec_detail = rec.get("detail", "")
            else:
                rec_title = str(rec)
                rec_detail = ""

            num_label = str(idx + 1)

            # Number circle (LG RED)
            circle_size = Cm(1.2)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Cm(1.5), item_y, int(circle_size), int(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors.RED
            circle.line.fill.background()
            tf = circle.text_frame
            tf.clear()
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = num_label
            self._set_font(run, size=Pt(14), bold=True, color=self.colors.WHITE)

            # Recommendation title text
            text_left = Cm(3.2)
            self._add_textbox(
                slide, text_left, item_y,
                Cm(29), Cm(1.2),
                text=rec_title,
                size=Pt(16), bold=True,
                color=self.colors.BLACK,
                alignment=PP_ALIGN.LEFT,
                vertical=MSO_ANCHOR.MIDDLE
            )

            # Detail text (if provided)
            if rec_detail:
                self._add_textbox(
                    slide, text_left, item_y + int(Cm(1.2)),
                    Cm(29), Cm(1.0),
                    text=rec_detail,
                    size=LGTypography.BODY,
                    color=self.colors.DARK_GRAY,
                    alignment=PP_ALIGN.LEFT,
                    vertical=MSO_ANCHOR.TOP
                )
                item_y += int(Cm(2.5))
            else:
                item_y += int(Cm(1.8))

        return slide

    def add_closing(self, text="감사합니다", subtitle=""):
        """
        Add a closing/thank you slide.

        Args:
            text: Main closing text
            subtitle: Additional text
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # Theme-specific closing decoration
        if self.theme == "hanwha":
            # Hanwha: bottom orange accent bar
            bar_h = Cm(0.15)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Cm(0), LGDimensions.SLIDE_HEIGHT - bar_h,
                LGDimensions.SLIDE_WIDTH, bar_h
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.colors.RED
            bar.line.fill.background()
        else:
            # LG: L-brackets (same as cover)
            self._add_l_bracket(slide, "top-left")
            self._add_l_bracket(slide, "bottom-right")

        # Closing text (centered)
        self._add_textbox(
            slide,
            left=Cm(3), top=sh // 2 - Cm(2),
            width=sw - Cm(6), height=Cm(3),
            text=text,
            size=Pt(36),
            bold=True,
            color=self.colors.BLACK,
            alignment=PP_ALIGN.CENTER,
            vertical=MSO_ANCHOR.MIDDLE
        )

        if subtitle:
            self._add_textbox(
                slide,
                left=Cm(3), top=sh // 2 + Cm(1.5),
                width=sw - Cm(6), height=Cm(1.5),
                text=subtitle,
                size=LGTypography.COVER_SUBTITLE,
                color=self.colors.DARK_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        # Logo
        if self.logo_path and os.path.exists(self.logo_path):
            slide.shapes.add_picture(
                self.logo_path,
                sw - Cm(5), sh - Cm(4),
                width=Cm(3)
            )

        return slide

    # ─────────────────────────────────────────
    # Summary Matrix / Pitch Deck Templates
    # ─────────────────────────────────────────

    def add_summary_matrix(self, title, section="", headers=None,
                           row_groups=None, subtitle="", chapter=""):
        """
        Add a summary matrix table — grouped rows with merged category cells.
        Matches the LG pattern of category labels on the left with detailed
        content cells across columns.

        Args:
            title: Slide title (e.g., "공정 DX팀 시스템 현황 Summary")
            section: Section name for top-right indicator
            headers: List of column header strings (e.g., ["LG에너지솔루션", "LG디스플레이", ...])
            row_groups: List of dicts, each representing a merged group:
                [
                    {
                        "category": "공정제어",      # Left-most merged label
                        "rows": [
                            {
                                "sub_label": "과제\\n현황",  # Sub-category label
                                "cells": ["셀1 내용", "셀2 내용", ...]  # One per header
                            },
                            {
                                "sub_label": "AI 적용\\n과제",
                                "cells": ["셀1 내용", "셀2 내용", ...]
                            }
                        ]
                    }
                ]
            subtitle: Optional subtitle text
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_start = int(y_pos)

        if not headers or not row_groups:
            return slide

        # Calculate table dimensions
        # Columns: category_col + sub_label_col + data_cols
        num_data_cols = len(headers)
        num_cols = 2 + num_data_cols  # category + sub_label + data columns
        total_data_rows = sum(len(g["rows"]) for g in row_groups)
        num_rows = 1 + total_data_rows  # header + data

        table_width = sw - int(Cm(2.0))
        available_h = int(LGDimensions.SLIDE_HEIGHT) - y_start - int(Cm(0.8))
        row_h = min(Cm(2.5), available_h // num_rows)
        table_height = row_h * num_rows

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Cm(1.5), y_start, table_width, table_height
        )
        table = table_shape.table

        # Disable default styling
        table.first_row = False
        table.first_col = False
        table.last_row = False
        table.last_col = False
        table.horz_banding = False
        table.vert_banding = False

        # Column widths
        cat_width = int(table_width * 0.07)
        sub_width = int(table_width * 0.07)
        data_width = (table_width - cat_width - sub_width) // num_data_cols
        table.columns[0].width = cat_width
        table.columns[1].width = sub_width
        for c in range(num_data_cols):
            table.columns[2 + c].width = data_width

        # Header row: empty for first 2 cols, then data headers
        self._format_table_cell(table.cell(0, 0), "", is_header=False,
                                bg_color=self.colors.WHITE)
        self._set_cell_border(table.cell(0, 0), color="CCCCCC")
        self._format_table_cell(table.cell(0, 1), "", is_header=False,
                                bg_color=self.colors.WHITE)
        self._set_cell_border(table.cell(0, 1), color="CCCCCC")
        for c, header in enumerate(headers):
            self._format_table_cell(table.cell(0, 2 + c), header, is_header=True)
            self._set_cell_border(table.cell(0, 2 + c), color="3C3C3C")

        # Data rows with merged category cells
        current_row = 1
        for group in row_groups:
            category = group["category"]
            group_rows = group["rows"]
            group_start = current_row
            group_end = current_row + len(group_rows) - 1

            for i, row_data in enumerate(group_rows):
                sub_label = row_data.get("sub_label", "")
                cells = row_data.get("cells", [])
                r = current_row

                # Sub-label cell (gray background)
                self._format_table_cell(
                    table.cell(r, 1), sub_label.replace("\\n", "\n"),
                    font_size=Pt(9), bold=True, bg_color=self.colors.LIGHT_GRAY,
                    alignment=PP_ALIGN.CENTER
                )
                self._set_cell_border(table.cell(r, 1), color="CCCCCC")

                # Data cells
                for c, cell_text in enumerate(cells):
                    if c < num_data_cols:
                        self._format_table_cell(
                            table.cell(r, 2 + c), cell_text,
                            font_size=Pt(9), alignment=PP_ALIGN.LEFT,
                            bg_color=self.colors.WHITE
                        )
                        self._set_cell_border(table.cell(r, 2 + c), color="CCCCCC")

                current_row += 1

            # Merge category cells vertically
            if len(group_rows) > 1:
                table.cell(group_start, 0).merge(table.cell(group_end, 0))
            self._format_table_cell(
                table.cell(group_start, 0), category,
                font_size=Pt(10), bold=True, bg_color=self.colors.LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER
            )
            self._set_cell_border(table.cell(group_start, 0), color="CCCCCC")

        return slide

    def add_two_column(self, title, section="", left_title="", left_bullets=None,
                       right_title="", right_bullets=None, subtitle="",
                       chapter=""):
        """
        Add a two-column layout slide — useful for comparisons, pros/cons,
        before/after, or side-by-side content.

        Args:
            title: Slide title
            section: Section name for top-right
            left_title: Title for left column
            left_bullets: List of bullet strings for left column
            right_title: Title for right column
            right_bullets: List of bullet strings for right column
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        col_width = (sw - int(Cm(3.5))) // 2
        left_x = Cm(1.5)
        right_x = int(left_x) + col_width + int(Cm(0.5))
        header_h = Cm(1.0)
        body_h = int(LGDimensions.SLIDE_HEIGHT) - y_pos - int(header_h) - int(Cm(1.5))

        for col_x, col_title, col_bullets in [
            (left_x, left_title, left_bullets or []),
            (right_x, right_title, right_bullets or [])
        ]:
            # Column header box (DARK_RED background with white text)
            if col_title:
                hdr_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    int(col_x), y_pos, col_width, int(header_h)
                )
                hdr_box.fill.solid()
                hdr_box.fill.fore_color.rgb = self.colors.DARK_RED
                hdr_box.line.fill.background()
                tf = hdr_box.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Cm(0.4)
                tf.margin_right = Cm(0.3)
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = col_title
                self._set_font(run, size=Pt(14), bold=True,
                              color=self.colors.WHITE)

            # Column body box (WHITE background with bullets, rounded)
            body_top = y_pos + int(header_h)
            body_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(col_x), body_top, col_width, body_h
            )
            body_box.fill.solid()
            body_box.fill.fore_color.rgb = self.colors.WHITE
            body_box.line.color.rgb = self.colors.BORDER_GRAY
            body_box.line.width = Pt(0.5)

            tf = body_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Cm(0.4)
            tf.margin_right = Cm(0.3)
            tf.margin_top = Cm(0.3)

            if col_bullets:
                for j, bullet in enumerate(col_bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(3)
                    pPr = p._p.get_or_add_pPr()
                    buChar = _make_oxml_element('a:buChar')
                    buChar.set('char', '\u2022')
                    pPr.append(buChar)
                    pPr.set('marL', str(int(Cm(0.5))))
                    pPr.set('indent', str(int(Cm(-0.4))))
                    run = p.add_run()
                    run.text = bullet
                    self._set_font(run, size=LGTypography.BODY,
                                  color=self.colors.BLACK)

            self._add_shadow(body_box)

        return slide

    # ─────────────────────────────────────────
    # Template-Based Layout Methods
    # ─────────────────────────────────────────

    def _density_sizes(self, density):
        """Return body/bullet font sizes for a given density level."""
        density_map = {
            "spacious": {"body": LGTypography.BODY_SPACIOUS, "bullet": LGTypography.BULLET_SPACIOUS},
            "normal":   {"body": LGTypography.BODY, "bullet": LGTypography.BODY},
            "compact":  {"body": LGTypography.BODY_COMPACT, "bullet": LGTypography.BULLET_COMPACT},
            "dense":    {"body": LGTypography.BODY_DENSE, "bullet": LGTypography.BULLET_DENSE},
        }
        return density_map.get(density, density_map["normal"])

    def _add_title_pill(self, slide, x, y, text, width=None, height=None,
                        bg_color=None):
        """
        Add a DARK_RED title pill badge centered inside a content box.
        Returns the shape.
        """
        w = width or Cm(5.17)
        h = height or Cm(0.85)
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = bg_color or self.colors.DARK_RED
        pill.line.fill.background()
        tf = pill.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Cm(0.15)
        tf.margin_right = Cm(0.15)
        tf.margin_top = Cm(0.05)
        tf.margin_bottom = Cm(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        self._set_font(run, size=Pt(10), bold=True, color=self.colors.WHITE)
        return pill

    def _add_circle_icon(self, slide, x, y, icon_text, size=None,
                         bg_color=None, text_color=None):
        """
        Add a circle icon shape with centered text.
        Default: LG RED background with white text.
        Returns the shape.
        """
        s = size or Cm(2.95)
        bg = bg_color or self.colors.RED
        txt_clr = text_color or self.colors.WHITE
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(x), int(y), int(s), int(s)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = bg
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon_text or ""
        self._set_font(run, size=Pt(24), bold=True, color=txt_clr)
        return circle

    def _add_bullets_to_frame(self, tf, bullets, font_size=None, start_index=0):
        """Add bullet items to an existing text frame."""
        sz = font_size or LGTypography.BODY
        for j, bullet in enumerate(bullets):
            idx = start_index + j
            p = tf.paragraphs[0] if idx == 0 and tf.paragraphs[0].text == '' else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(3)
            pPr = p._p.get_or_add_pPr()
            buChar = _make_oxml_element('a:buChar')
            buChar.set('char', '\u2022')
            pPr.append(buChar)
            pPr.set('marL', str(int(Cm(0.5))))
            pPr.set('indent', str(int(Cm(-0.4))))
            run = p.add_run()
            run.text = bullet
            self._set_font(run, size=sz, color=self.colors.BLACK)

    def _add_item_pill(self, slide, x, y, text, width=None, height=None):
        """Add a small rounded item pill (LIGHT_GRAY bg, BORDER_GRAY border)."""
        w = width or Cm(6.37)
        h = height or Cm(0.85)
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = self.colors.LIGHT_GRAY
        pill.line.color.rgb = self.colors.BORDER_GRAY
        pill.line.width = Pt(1)
        tf = pill.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Cm(0.2)
        tf.margin_right = Cm(0.2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        self._set_font(run, size=Pt(10), color=self.colors.BLACK)
        return pill

    # ── 1. Multi-Column Layout ──

    def add_multi_column(self, title, section="", columns=None, ratio=None,
                         subtitle="", chapter="", density="normal"):
        """
        Flexible multi-column layout with rounded content boxes.
        Replaces/supplements existing two_column for N-column layouts.

        Args:
            title: Slide title
            section: Section name for top-right tag
            columns: list of dicts, each with:
                - title (str): column header (shown as pill badge)
                - bullets (list[str]): bullet items
                - body (str): plain text body (alternative to bullets)
            ratio: column width ratios, e.g. [7,3] or [3,7] or [1,1,1] or [1,1,1,1]
                   defaults to equal widths
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
            density: "spacious" | "normal" | "compact" | "dense"
        """
        slide = self._get_blank_slide()
        columns = columns or []
        n_cols = len(columns) if columns else 1
        sizes = self._density_sizes(density)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        # Content area dimensions
        content_left = int(Cm(1.33))
        content_width = int(Cm(31.16))
        content_height = int(Cm(11.22))
        gap = int(Cm(0.35))

        # Calculate column widths from ratio
        if ratio and len(ratio) == n_cols:
            total_ratio = sum(ratio)
            available = content_width - gap * (n_cols - 1)
            col_widths = [int(available * r / total_ratio) for r in ratio]
        else:
            available = content_width - gap * (n_cols - 1)
            col_widths = [int(available / n_cols)] * n_cols

        # Draw columns
        x_cursor = content_left
        for i, col in enumerate(columns):
            col_w = col_widths[i] if i < len(col_widths) else col_widths[-1]
            col_title = col.get("title", "")
            col_bullets = col.get("bullets", [])
            col_body = col.get("body", "")

            # Rounded content box — subtle fill + shadow for depth (no border)
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_cursor, y_pos, col_w, content_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.colors.LIGHT_GRAY
            box.line.fill.background()
            # Set corner radius (template: val 8000)
            _spPr = box._element.spPr
            _pg = _spPr.find(qn('a:prstGeom'))
            if _pg is not None:
                _al = _pg.find(qn('a:avLst'))
                if _al is None:
                    _al = _make_oxml_element('a:avLst')
                    _pg.append(_al)
                for _gd in _al.findall(qn('a:gd')):
                    _al.remove(_gd)
                _gd = _make_oxml_element('a:gd')
                _gd.set('name', 'adj')
                _gd.set('fmla', 'val 8000')
                _al.append(_gd)

            self._add_shadow(box, shadow_type="card")

            # Title pill inside box (centered horizontally)
            inner_y = y_pos + int(Cm(0.70))
            if col_title:
                pill_w = min(int(Cm(5.17)), col_w - int(Cm(0.60)))
                pill_x = x_cursor + (col_w - pill_w) // 2
                self._add_title_pill(slide, pill_x, inner_y, col_title,
                                     width=pill_w)
                text_top = inner_y + int(Cm(1.10))
            else:
                text_top = inner_y

            # Content: bullets or body text
            text_h = content_height - (text_top - y_pos) - int(Cm(0.40))
            text_w = col_w - int(Cm(0.80))
            text_x = x_cursor + int(Cm(0.40))

            if col_bullets:
                tb = slide.shapes.add_textbox(text_x, text_top, text_w, text_h)
                tf = tb.text_frame
                tf.word_wrap = True
                tb.fill.background()
                tb.line.fill.background()
                self._add_bullets_to_frame(tf, col_bullets,
                                           font_size=sizes["bullet"])
            elif col_body:
                self._add_textbox(
                    slide, text_x, text_top, text_w, text_h,
                    text=col_body, size=sizes["body"], color=self.colors.BLACK
                )

            x_cursor += col_w + gap

        return slide

    # ── 2. Card Grid Layout ──

    def add_card_grid(self, title, section="", cards=None, cols=2,
                      subtitle="", chapter="", show_icon=True,
                      density="normal"):
        """
        Grid of titled card panels with optional icon circles.

        Args:
            title: Slide title
            section: Section name for top-right tag
            cards: list of dicts, each with:
                - title (str): card header (pill badge)
                - items (list[str]): list items inside card
                - icon_text (str): text inside circle icon (1-2 chars)
                - footer (str): footer text at bottom of card
            cols: number of columns (2, 3, or 4)
            show_icon: whether to show circle icons (default True)
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
            density: "spacious" | "normal" | "compact" | "dense"
        """
        slide = self._get_blank_slide()
        cards = cards or []
        sizes = self._density_sizes(density)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        # Content area
        content_left = int(Cm(1.33))
        content_width = int(Cm(31.16))
        content_height = int(Cm(11.22))
        gap = int(Cm(0.50))

        # Column widths
        cols = max(1, min(cols, 6))
        col_w = (content_width - gap * (cols - 1)) // cols

        # Arrange cards into rows
        for idx, card in enumerate(cards):
            col_idx = idx % cols
            row_idx = idx // cols
            row_h = content_height  # single row for simplicity
            if len(cards) > cols:
                row_h = (content_height - gap * ((len(cards) - 1) // cols)) // (((len(cards) - 1) // cols) + 1)

            card_x = content_left + col_idx * (col_w + gap)
            card_y = y_pos + row_idx * (row_h + gap)
            card_h = min(row_h, content_height - (card_y - y_pos))

            card_title = card.get("title", "")
            card_items = card.get("items", [])
            card_icon = card.get("icon_text", "")
            card_footer = card.get("footer", "")

            # Card background box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                card_x, card_y, col_w, card_h
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.colors.WHITE
            box.line.color.rgb = self.colors.RED
            box.line.width = Pt(1)
            # Set corner radius (template: val 8000)
            _spPr = box._element.spPr
            _pg = _spPr.find(qn('a:prstGeom'))
            if _pg is not None:
                _al = _pg.find(qn('a:avLst'))
                if _al is None:
                    _al = _make_oxml_element('a:avLst')
                    _pg.append(_al)
                for _gd in _al.findall(qn('a:gd')):
                    _al.remove(_gd)
                _gd = _make_oxml_element('a:gd')
                _gd.set('name', 'adj')
                _gd.set('fmla', 'val 8000')
                _al.append(_gd)

            self._add_shadow(box)

            inner_y = card_y + int(Cm(0.70))

            # Title pill
            if card_title:
                pill_w = min(int(Cm(5.17)), col_w - int(Cm(0.60)))
                pill_x = card_x + (col_w - pill_w) // 2
                self._add_title_pill(slide, pill_x, inner_y, card_title,
                                     width=pill_w)
                inner_y += int(Cm(1.10))

            # Circle icon
            if show_icon and card_icon:
                icon_size = int(Cm(2.95))
                icon_x = card_x + (col_w - icon_size) // 2
                self._add_circle_icon(slide, icon_x, inner_y, card_icon)
                inner_y += icon_size + int(Cm(0.40))

            # Item pills (stacked)
            if card_items:
                item_w = col_w - int(Cm(1.00))
                item_x = card_x + (col_w - item_w) // 2
                for item_text in card_items:
                    if inner_y + int(Cm(0.85)) > card_y + card_h - int(Cm(0.80)):
                        break  # prevent overflow
                    self._add_item_pill(slide, item_x, inner_y, item_text,
                                        width=item_w)
                    inner_y += int(Cm(1.05))

            # Footer text
            if card_footer:
                footer_y = card_y + card_h - int(Cm(0.90))
                footer_w = col_w - int(Cm(0.60))
                footer_x = card_x + int(Cm(0.30))
                self._add_textbox(
                    slide, footer_x, footer_y, footer_w, Cm(0.70),
                    text=card_footer, size=Pt(9), bold=True,
                    color=self.colors.RED, alignment=PP_ALIGN.CENTER
                )

        return slide

    # ── 3. Content Sidebar Layout ──

    def add_content_sidebar(self, title, section="", main_body="",
                            main_bullets=None, sidebar_items=None,
                            subtitle="", chapter="", sidebar_position="right",
                            sidebar_ratio=0.48, density="normal"):
        """
        Main content area with stacked sidebar panels.

        Args:
            title: Slide title
            section: Section name for top-right tag
            main_body: Body text for main area
            main_bullets: Bullet list for main area
            sidebar_items: list of dicts for sidebar panels:
                - title (str): panel header
                - body (str): panel content text
                - bullets (list[str]): panel bullet items
            sidebar_position: "right" or "left"
            sidebar_ratio: sidebar width as fraction of total (default 0.48)
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
            density: "spacious" | "normal" | "compact" | "dense"
        """
        slide = self._get_blank_slide()
        sidebar_items = sidebar_items or []
        sizes = self._density_sizes(density)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        # Content area
        content_left = int(Cm(1.33))
        content_width = int(Cm(31.16))
        content_height = int(Cm(11.22))
        gap = int(Cm(0.50))

        sidebar_w = int(content_width * sidebar_ratio)
        main_w = content_width - sidebar_w - gap

        if sidebar_position == "left":
            sidebar_x = content_left
            main_x = content_left + sidebar_w + gap
        else:
            main_x = content_left
            sidebar_x = content_left + main_w + gap

        # Main content box
        main_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            main_x, y_pos, main_w, content_height
        )
        main_box.fill.solid()
        main_box.fill.fore_color.rgb = self.colors.WHITE
        main_box.line.color.rgb = self.colors.RED
        main_box.line.width = Pt(1)
        # Set corner radius (template: val 8000)
        _spPr = main_box._element.spPr
        _pg = _spPr.find(qn('a:prstGeom'))
        if _pg is not None:
            _al = _pg.find(qn('a:avLst'))
            if _al is None:
                _al = _make_oxml_element('a:avLst')
                _pg.append(_al)
            for _gd in _al.findall(qn('a:gd')):
                _al.remove(_gd)
            _gd = _make_oxml_element('a:gd')
            _gd.set('name', 'adj')
            _gd.set('fmla', 'val 8000')
            _al.append(_gd)

        self._add_shadow(main_box)

        # Main content text
        text_x = main_x + int(Cm(0.50))
        text_y = y_pos + int(Cm(0.50))
        text_w = main_w - int(Cm(1.00))
        text_h = content_height - int(Cm(1.00))

        if main_bullets:
            tb = slide.shapes.add_textbox(text_x, text_y, text_w, text_h)
            tf = tb.text_frame
            tf.word_wrap = True
            tb.fill.background()
            tb.line.fill.background()
            self._add_bullets_to_frame(tf, main_bullets,
                                       font_size=sizes["bullet"])
        elif main_body:
            self._add_textbox(
                slide, text_x, text_y, text_w, text_h,
                text=main_body, size=sizes["body"], color=self.colors.BLACK
            )

        # Sidebar panels (stacked vertically)
        n_panels = max(len(sidebar_items), 1)
        panel_gap = int(Cm(0.50))
        total_gaps = panel_gap * (n_panels - 1)
        panel_h = (content_height - total_gaps) // n_panels

        panel_y = y_pos
        for panel in sidebar_items:
            panel_title = panel.get("title", "")
            panel_body = panel.get("body", "")
            panel_bullets = panel.get("bullets", [])

            # Panel box
            pbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                sidebar_x, panel_y, sidebar_w, panel_h
            )
            pbox.fill.solid()
            pbox.fill.fore_color.rgb = self.colors.WHITE
            pbox.line.color.rgb = self.colors.RED
            pbox.line.width = Pt(1)
            # Set corner radius (template: val 8000)
            _spPr = pbox._element.spPr
            _pg = _spPr.find(qn('a:prstGeom'))
            if _pg is not None:
                _al = _pg.find(qn('a:avLst'))
                if _al is None:
                    _al = _make_oxml_element('a:avLst')
                    _pg.append(_al)
                for _gd in _al.findall(qn('a:gd')):
                    _al.remove(_gd)
                _gd = _make_oxml_element('a:gd')
                _gd.set('name', 'adj')
                _gd.set('fmla', 'val 8000')
                _al.append(_gd)

            self._add_shadow(pbox)

            inner_y = panel_y + int(Cm(0.40))
            inner_x = sidebar_x + int(Cm(0.40))
            inner_w = sidebar_w - int(Cm(0.80))

            # Panel title (bold text at top)
            if panel_title:
                self._add_textbox(
                    slide, inner_x, inner_y, inner_w, Cm(0.70),
                    text=panel_title, size=Pt(12), bold=True,
                    color=self.colors.BLACK
                )
                inner_y += int(Cm(0.80))

            # Panel body or bullets
            remaining_h = panel_h - (inner_y - panel_y) - int(Cm(0.20))
            if panel_bullets:
                tb = slide.shapes.add_textbox(inner_x, inner_y,
                                              inner_w, remaining_h)
                tf = tb.text_frame
                tf.word_wrap = True
                tb.fill.background()
                tb.line.fill.background()
                self._add_bullets_to_frame(tf, panel_bullets,
                                           font_size=sizes["bullet"])
            elif panel_body:
                self._add_textbox(
                    slide, inner_x, inner_y, inner_w, remaining_h,
                    text=panel_body, size=sizes["body"],
                    color=self.colors.BLACK
                )

            panel_y += panel_h + panel_gap

        return slide

    # ── 4. Titled Columns Layout ──

    def add_titled_columns(self, title, section="", bar_title="",
                           columns=None, subtitle="", chapter="",
                           show_bar=True, density="normal"):
        """
        Full-width title bar with columns below.

        Args:
            title: Slide title
            section: Section name for top-right tag
            bar_title: text for the full-width title bar
            columns: list of dicts:
                - title (str): column header (pill badge)
                - icon_text (str): text in circle icon
                - bullets (list[str]): bullet items
                - items (list[str]): bottom item pills
            show_bar: whether to show the full-width title bar (default True)
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
            density: "spacious" | "normal" | "compact" | "dense"
        """
        slide = self._get_blank_slide()
        columns = columns or []
        n_cols = len(columns) if columns else 1
        sizes = self._density_sizes(density)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        # Content area
        content_left = int(Cm(1.33))
        content_width = int(Cm(31.16))
        bar_height = int(Cm(1.07))
        gap = int(Cm(0.40))

        # Title bar
        if show_bar and bar_title:
            bar_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                content_left, y_pos, content_width, bar_height
            )
            bar_box.fill.solid()
            bar_box.fill.fore_color.rgb = self.colors.WHITE
            bar_box.line.color.rgb = self.colors.RED
            bar_box.line.width = Pt(1)
            # Set corner radius (template: val 8000)
            _spPr = bar_box._element.spPr
            _pg = _spPr.find(qn('a:prstGeom'))
            if _pg is not None:
                _al = _pg.find(qn('a:avLst'))
                if _al is None:
                    _al = _make_oxml_element('a:avLst')
                    _pg.append(_al)
                for _gd in _al.findall(qn('a:gd')):
                    _al.remove(_gd)
                _gd = _make_oxml_element('a:gd')
                _gd.set('name', 'adj')
                _gd.set('fmla', 'val 8000')
                _al.append(_gd)

            self._add_shadow(bar_box)

            tf = bar_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Cm(0.4)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = bar_title
            self._set_font(run, size=Pt(12), bold=True, color=self.colors.BLACK)

            col_top = y_pos + bar_height + gap
        else:
            col_top = y_pos

        # Column area height
        content_bottom = y_pos + int(Cm(11.22))
        col_height = content_bottom - col_top
        col_gap = int(Cm(0.50))
        col_w = (content_width - col_gap * (n_cols - 1)) // n_cols

        # Draw columns
        x_cursor = content_left
        for col in columns:
            col_title = col.get("title", "")
            col_icon = col.get("icon_text", "")
            col_bullets = col.get("bullets", [])
            col_items = col.get("items", [])

            # Column background box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_cursor, col_top, col_w, col_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.colors.WHITE
            box.line.color.rgb = self.colors.RED
            box.line.width = Pt(1)
            # Set corner radius (template: val 8000)
            _spPr = box._element.spPr
            _pg = _spPr.find(qn('a:prstGeom'))
            if _pg is not None:
                _al = _pg.find(qn('a:avLst'))
                if _al is None:
                    _al = _make_oxml_element('a:avLst')
                    _pg.append(_al)
                for _gd in _al.findall(qn('a:gd')):
                    _al.remove(_gd)
                _gd = _make_oxml_element('a:gd')
                _gd.set('name', 'adj')
                _gd.set('fmla', 'val 8000')
                _al.append(_gd)

            self._add_shadow(box)

            inner_y = col_top + int(Cm(0.70))

            # Title pill
            if col_title:
                pill_w = min(int(Cm(5.17)), col_w - int(Cm(0.60)))
                pill_x = x_cursor + (col_w - pill_w) // 2
                self._add_title_pill(slide, pill_x, inner_y, col_title,
                                     width=pill_w)
                inner_y += int(Cm(1.10))

            # Circle icon
            if col_icon:
                icon_size = int(Cm(2.95))
                icon_x = x_cursor + (col_w - icon_size) // 2
                self._add_circle_icon(slide, icon_x, inner_y, col_icon)
                inner_y += icon_size + int(Cm(0.40))

            # Bullets
            if col_bullets:
                bullet_h = col_height - (inner_y - col_top) - int(Cm(0.50))
                if col_items:
                    bullet_h -= int(Cm(1.10)) * len(col_items)
                bullet_w = col_w - int(Cm(0.80))
                bullet_x = x_cursor + int(Cm(0.40))
                tb = slide.shapes.add_textbox(bullet_x, inner_y,
                                              bullet_w, max(bullet_h, int(Cm(1.0))))
                tf = tb.text_frame
                tf.word_wrap = True
                tb.fill.background()
                tb.line.fill.background()
                self._add_bullets_to_frame(tf, col_bullets,
                                           font_size=sizes["bullet"])
                inner_y += max(bullet_h, int(Cm(1.0)))

            # Bottom item pills
            if col_items:
                item_w = col_w - int(Cm(1.00))
                item_x = x_cursor + (col_w - item_w) // 2
                # Position items from the bottom of the column
                items_start = col_top + col_height - int(Cm(0.50)) - int(Cm(1.05)) * len(col_items)
                pill_y = max(inner_y, items_start)
                for item_text in col_items:
                    if pill_y + int(Cm(0.85)) > col_top + col_height:
                        break
                    self._add_item_pill(slide, item_x, pill_y, item_text,
                                        width=item_w)
                    pill_y += int(Cm(1.05))

            x_cursor += col_w + col_gap

        return slide

    # ── 5. Icon Process Layout ──

    def add_icon_process(self, title, section="", steps=None,
                         bottom_sections=None, subtitle="", chapter="",
                         density="normal"):
        """
        Process flow with icon groups on top and detail sections below.

        Args:
            title: Slide title
            section: Section name for top-right tag
            steps: list of dicts for top process:
                - title (str): step name (bold)
                - description (str): step description
                - icon_text (str): icon circle text (1-2 chars)
            bottom_sections: list of dicts for bottom detail panels:
                - title (str): section header (pill badge)
                - items (list of dicts with 'label' and 'value')
                - bullets (list[str]): alternative to items
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
            density: "spacious" | "normal" | "compact" | "dense"
        """
        slide = self._get_blank_slide()
        steps = steps or []
        bottom_sections = bottom_sections or []
        sizes = self._density_sizes(density)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        # Content area
        content_left = int(Cm(1.33))
        content_width = int(Cm(31.16))
        content_height = int(Cm(11.22))

        # Top process area: ~6.5cm, bottom detail area: ~4.0cm
        top_height = int(Cm(6.50))
        bottom_gap = int(Cm(0.50))
        bottom_top = y_pos + top_height + bottom_gap
        bottom_height = content_height - top_height - bottom_gap

        # ── Top: Icon steps with arrows ──
        n_steps = max(len(steps), 1)
        icon_size = int(Cm(2.95))
        arrow_w = int(Cm(1.20))

        # Space available for steps (icons + arrows between)
        total_arrows = max(n_steps - 1, 0)
        total_arrow_space = total_arrows * arrow_w
        step_space = (content_width - total_arrow_space) // n_steps

        for i, step in enumerate(steps):
            step_title = step.get("title", "")
            step_desc = step.get("description", "")
            step_icon = step.get("icon_text", "")

            step_x = content_left + i * (step_space + arrow_w)
            step_center = step_x + step_space // 2

            # Circle icon (centered in step area)
            icon_x = step_center - icon_size // 2
            icon_y = y_pos + int(Cm(0.30))
            self._add_circle_icon(slide, icon_x, icon_y, step_icon)

            # Title below icon
            title_y = icon_y + icon_size + int(Cm(0.30))
            self._add_textbox(
                slide, step_x, title_y, step_space, Cm(0.70),
                text=step_title, size=Pt(11), bold=True,
                color=self.colors.BLACK, alignment=PP_ALIGN.CENTER
            )

            # Description below title
            if step_desc:
                desc_y = title_y + int(Cm(0.80))
                self._add_textbox(
                    slide, step_x, desc_y, step_space, Cm(1.50),
                    text=step_desc, size=sizes["body"],
                    color=self.colors.DARK_GRAY, alignment=PP_ALIGN.CENTER
                )

            # Arrow between steps (chevron triangle)
            if i < n_steps - 1:
                arrow_x = step_x + step_space
                arrow_y = icon_y + icon_size // 2 - int(Cm(0.40))
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.CHEVRON,
                    arrow_x, arrow_y, arrow_w, int(Cm(0.80))
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors.RED
                arrow.line.fill.background()

        # ── Bottom: Detail panels ──
        if bottom_sections:
            n_panels = len(bottom_sections)
            panel_gap = int(Cm(0.50))
            panel_w = (content_width - panel_gap * (n_panels - 1)) // n_panels

            panel_x = content_left
            for panel in bottom_sections:
                panel_title = panel.get("title", "")
                panel_items = panel.get("items", [])
                panel_bullets = panel.get("bullets", [])

                # Panel background box
                pbox = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    panel_x, bottom_top, panel_w, bottom_height
                )
                pbox.fill.solid()
                pbox.fill.fore_color.rgb = self.colors.WHITE
                pbox.line.color.rgb = self.colors.RED
                pbox.line.width = Pt(1)
                # Set corner radius (template: val 8000)
                _spPr = pbox._element.spPr
                _pg = _spPr.find(qn('a:prstGeom'))
                if _pg is not None:
                    _al = _pg.find(qn('a:avLst'))
                    if _al is None:
                        _al = _make_oxml_element('a:avLst')
                        _pg.append(_al)
                    for _gd in _al.findall(qn('a:gd')):
                        _al.remove(_gd)
                    _gd = _make_oxml_element('a:gd')
                    _gd.set('name', 'adj')
                    _gd.set('fmla', 'val 8000')
                    _al.append(_gd)

                self._add_shadow(pbox)

                inner_y = bottom_top + int(Cm(0.50))

                # Panel title pill
                if panel_title:
                    pill_w = min(int(Cm(5.17)), panel_w - int(Cm(0.60)))
                    pill_x = panel_x + (panel_w - pill_w) // 2
                    self._add_title_pill(slide, pill_x, inner_y, panel_title,
                                         width=pill_w)
                    inner_y += int(Cm(1.10))

                # Items (label: value pairs)
                if panel_items:
                    item_x = panel_x + int(Cm(0.50))
                    item_w = panel_w - int(Cm(1.00))
                    for item in panel_items:
                        label = item.get("label", "")
                        value = item.get("value", "")
                        display = f"{label}: {value}" if label and value else label or value
                        self._add_textbox(
                            slide, item_x, inner_y, item_w, Cm(0.60),
                            text=display, size=sizes["body"],
                            color=self.colors.BLACK
                        )
                        inner_y += int(Cm(0.65))

                # Bullets alternative
                elif panel_bullets:
                    remaining_h = bottom_height - (inner_y - bottom_top) - int(Cm(0.20))
                    tb = slide.shapes.add_textbox(
                        panel_x + int(Cm(0.50)), inner_y,
                        panel_w - int(Cm(1.00)), remaining_h
                    )
                    tf = tb.text_frame
                    tf.word_wrap = True
                    tb.fill.background()
                    tb.line.fill.background()
                    self._add_bullets_to_frame(tf, panel_bullets,
                                               font_size=sizes["bullet"])

                panel_x += panel_w + panel_gap

        return slide

    def add_kpi_cards(self, title, section="", kpis=None, subtitle="",
                      chapter=""):
        """
        Add a KPI/metrics card slide — large numbers with labels,
        commonly used in executive summaries and dashboards.

        Args:
            title: Slide title
            section: Section name for top-right
            kpis: List of dicts with KPI data:
                [
                    {"value": "30%", "label": "불량률 감소", "color": "#2E7D32"},
                    {"value": "15%", "label": "공정 효율 향상"},
                    {"value": "2.5억", "label": "연간 비용 절감", "color": "#1565C0"},
                ]
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        if not kpis:
            return slide

        y_pos += int(Cm(0.3))

        num_cards = len(kpis)
        card_gap = Cm(0.3)
        total_gap = int(card_gap) * (num_cards - 1)
        card_width = (sw - int(Cm(2.0)) - total_gap) // num_cards
        card_height = Cm(8)
        x_start = Cm(1.5)

        for i, kpi in enumerate(kpis):
            x = int(x_start) + i * (card_width + int(card_gap))

            # Card background (WHITE rounded rectangle)
            card = self.add_box(
                slide, x, y_pos, card_width, card_height,
                bg_color=self.colors.WHITE, rounded=True
            )

            # Value (large number)
            value_color_hex = kpi.get("color", None)
            if value_color_hex:
                r, g, b = int(value_color_hex[1:3], 16), int(value_color_hex[3:5], 16), int(value_color_hex[5:7], 16)
                v_color = RGBColor(r, g, b)
            else:
                v_color = self.colors.RED

            self._add_textbox(
                slide, x + Cm(0.5), y_pos + int(Cm(1.5)),
                card_width - Cm(1.0), Cm(3.5),
                text=kpi.get("value", ""),
                size=Pt(40), bold=True, color=v_color,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

            # Label
            self._add_textbox(
                slide, x + Cm(0.5), y_pos + int(Cm(5.5)),
                card_width - Cm(1.0), Cm(2.0),
                text=kpi.get("label", ""),
                size=LGTypography.BODY, color=self.colors.DARK_GRAY,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.TOP
            )

        return slide

    def add_timeline(self, title, section="", milestones=None, subtitle="",
                     chapter=""):
        """
        Add a horizontal timeline slide — useful for project phases,
        history, or sequential milestones.

        Args:
            title: Slide title
            section: Section name for top-right
            milestones: List of dicts:
                [
                    {"date": "2025 Q1", "title": "Phase 1", "description": "설명"},
                    {"date": "2025 Q2", "title": "Phase 2", "description": "설명"},
                ]
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        if not milestones:
            return slide

        num = len(milestones)
        line_y = y_pos + int(Cm(4.0))
        margin_x = Cm(1.5)
        line_width = sw - int(margin_x) * 2

        # Horizontal line
        self._add_horizontal_line(
            slide, margin_x, line_y, line_width,
            color=self.colors.CHARCOAL, height=Pt(3)
        )

        # Milestones
        spacing = line_width // num
        for i, ms in enumerate(milestones):
            cx = int(margin_x) + spacing * i + spacing // 2

            # Circle marker
            dot_size = Cm(0.6)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                cx - int(dot_size) // 2, line_y - int(dot_size) // 2,
                dot_size, dot_size
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.colors.RED
            dot.line.fill.background()

            # Date above line
            self._add_textbox(
                slide, cx - spacing // 2, line_y - Cm(2.5),
                spacing, Cm(2.0),
                text=ms.get("date", ""),
                size=Pt(11), bold=True, color=self.colors.RED,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.BOTTOM
            )

            # Title + description below line
            self._add_textbox(
                slide, cx - spacing // 2, line_y + Cm(1.0),
                spacing, Cm(1.5),
                text=ms.get("title", ""),
                size=LGTypography.BODY, bold=True, color=self.colors.BLACK,
                alignment=PP_ALIGN.CENTER
            )

            if ms.get("description"):
                self._add_textbox(
                    slide, cx - spacing // 2, line_y + Cm(2.5),
                    spacing, Cm(4.0),
                    text=ms["description"],
                    size=LGTypography.BODY_SMALL, color=self.colors.DARK_GRAY,
                    alignment=PP_ALIGN.CENTER
                )

        return slide

    def add_process_flow(self, title, section="", steps=None, subtitle="",
                         chapter=""):
        """
        Add a process flow slide — horizontal arrow-connected steps.
        Good for workflows, system architecture overview, or methodology.

        Args:
            title: Slide title
            section: Section name for top-right
            steps: List of dicts:
                [
                    {"title": "데이터 수집", "items": ["센서 데이터", "로그 수집"]},
                    {"title": "전처리", "items": ["정규화", "이상치 제거"]},
                    {"title": "AI 분석", "items": ["모델 학습", "추론"]},
                    {"title": "대시보드", "items": ["시각화", "알림"]},
                ]
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        if not steps:
            return slide

        num = len(steps)
        arrow_width = Cm(0.8)

        # Serpentine layout for >4 steps
        if num > 4:
            row1_count = (num + 1) // 2  # first row gets half (rounded up)
            row2_count = num - row1_count
            total_arrows_r1 = int(arrow_width) * (row1_count - 1)
            step_width = (sw - int(Cm(3.0)) - total_arrows_r1) // row1_count
            step_height = Cm(4.5)
            row_gap = Cm(1.2)
            x_start = Cm(1.5)

            # --- Row 1: left to right ---
            for i in range(row1_count):
                step = steps[i]
                box_x = int(x_start) + i * (step_width + int(arrow_width))

                header = self.add_box(
                    slide, box_x, y_pos, step_width, Cm(1.2),
                    text=step.get("title", ""),
                    bg_color=self.colors.DARK_RED, text_color=self.colors.WHITE,
                    font_size=LGTypography.BODY, bold=True, rounded=True
                )
                content = self.add_box(
                    slide, box_x, y_pos + int(Cm(1.2)),
                    step_width, step_height - Cm(1.2),
                    bg_color=self.colors.WHITE, rounded=True
                )

                items = step.get("items", [])
                desc = step.get("description", "")
                if not items and desc:
                    items = [line.strip() for line in desc.split("\n") if line.strip()]
                if items:
                    items_box = slide.shapes.add_textbox(
                        box_x + Cm(0.3), y_pos + int(Cm(1.6)),
                        step_width - Cm(0.6), step_height - Cm(2.0)
                    )
                    tf = items_box.text_frame
                    tf.word_wrap = True
                    items_box.fill.background()
                    items_box.line.fill.background()
                    for j, item in enumerate(items):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.space_after = Pt(3)
                        pPr = p._p.get_or_add_pPr()
                        buChar = _make_oxml_element('a:buChar')
                        buChar.set('char', '\u2022')
                        pPr.append(buChar)
                        run = p.add_run()
                        run.text = item
                        self._set_font(run, size=LGTypography.BODY_SMALL,
                                      color=self.colors.BLACK)

                if i < row1_count - 1:
                    arrow_x = box_x + step_width
                    arrow_y = y_pos + int(step_height) // 2 - int(Cm(0.5))
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        arrow_x, arrow_y, arrow_width, Cm(1.0)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = self.colors.RED
                    arrow.line.fill.background()

            # Down arrow from last in row1 to first in row2
            last_r1_x = int(x_start) + (row1_count - 1) * (step_width + int(arrow_width))
            down_arrow_x = last_r1_x + step_width // 2 - int(Cm(0.5))
            down_arrow_y = y_pos + int(step_height)
            down_arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                down_arrow_x, down_arrow_y, Cm(1.0), row_gap
            )
            down_arrow.fill.solid()
            down_arrow.fill.fore_color.rgb = self.colors.RED
            down_arrow.line.fill.background()

            # --- Row 2: left to right ---
            y_pos2 = y_pos + int(step_height) + int(row_gap)
            for i in range(row2_count):
                step = steps[row1_count + i]
                box_x = int(x_start) + i * (step_width + int(arrow_width))

                header = self.add_box(
                    slide, box_x, y_pos2, step_width, Cm(1.2),
                    text=step.get("title", ""),
                    bg_color=self.colors.DARK_RED, text_color=self.colors.WHITE,
                    font_size=LGTypography.BODY, bold=True, rounded=True
                )
                content = self.add_box(
                    slide, box_x, y_pos2 + int(Cm(1.2)),
                    step_width, step_height - Cm(1.2),
                    bg_color=self.colors.WHITE, rounded=True
                )

                items = step.get("items", [])
                desc = step.get("description", "")
                if not items and desc:
                    items = [line.strip() for line in desc.split("\n") if line.strip()]
                if items:
                    items_box = slide.shapes.add_textbox(
                        box_x + Cm(0.3), y_pos2 + int(Cm(1.6)),
                        step_width - Cm(0.6), step_height - Cm(2.0)
                    )
                    tf = items_box.text_frame
                    tf.word_wrap = True
                    items_box.fill.background()
                    items_box.line.fill.background()
                    for j, item in enumerate(items):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.space_after = Pt(3)
                        pPr = p._p.get_or_add_pPr()
                        buChar = _make_oxml_element('a:buChar')
                        buChar.set('char', '\u2022')
                        pPr.append(buChar)
                        run = p.add_run()
                        run.text = item
                        self._set_font(run, size=LGTypography.BODY_SMALL,
                                      color=self.colors.BLACK)

                if i < row2_count - 1:
                    arrow_x = box_x + step_width
                    arrow_y = y_pos2 + int(step_height) // 2 - int(Cm(0.5))
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        arrow_x, arrow_y, arrow_width, Cm(1.0)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = self.colors.RED
                    arrow.line.fill.background()

        else:
            # Original layout for <=4 steps
            total_arrows = int(arrow_width) * (num - 1)
            step_width = (sw - int(Cm(3.0)) - total_arrows) // num
            step_height = Cm(8)
            x = Cm(1.5)

            for i, step in enumerate(steps):
                box_x = int(x) + i * (step_width + int(arrow_width))

                header = self.add_box(
                    slide, box_x, y_pos, step_width, Cm(1.5),
                    text=step.get("title", ""),
                    bg_color=self.colors.DARK_RED, text_color=self.colors.WHITE,
                    font_size=LGTypography.BODY, bold=True, rounded=True
                )
                content = self.add_box(
                    slide, box_x, y_pos + int(Cm(1.5)),
                    step_width, step_height - Cm(1.5),
                    bg_color=self.colors.WHITE, rounded=True
                )

                items = step.get("items", [])
                desc = step.get("description", "")
                if not items and desc:
                    items = [line.strip() for line in desc.split("\n") if line.strip()]
                if items:
                    items_box = slide.shapes.add_textbox(
                        box_x + Cm(0.3), y_pos + int(Cm(2.2)),
                        step_width - Cm(0.6), step_height - Cm(3.0)
                    )
                    tf = items_box.text_frame
                    tf.word_wrap = True
                    items_box.fill.background()
                    items_box.line.fill.background()

                    for j, item in enumerate(items):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.space_after = Pt(3)
                        pPr = p._p.get_or_add_pPr()
                        buChar = _make_oxml_element('a:buChar')
                        buChar.set('char', '\u2022')
                        pPr.append(buChar)
                        run = p.add_run()
                        run.text = item
                        self._set_font(run, size=LGTypography.BODY_SMALL,
                                      color=self.colors.BLACK)

                if i < num - 1:
                    arrow_x = box_x + step_width
                    arrow_y = y_pos + int(step_height) // 2 - int(Cm(0.5))
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        arrow_x, arrow_y, arrow_width, Cm(1.0)
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = self.colors.RED
                    arrow.line.fill.background()

        return slide

    def add_swot(self, title, section="", strengths=None, weaknesses=None,
                 opportunities=None, threats=None, subtitle="", chapter=""):
        """
        Add a SWOT analysis slide — 2x2 grid with color-coded quadrants.

        Args:
            title: Slide title
            section: Section name for top-right
            strengths: List of strength strings
            weaknesses: List of weakness strings
            opportunities: List of opportunity strings
            threats: List of threat strings
            subtitle: Optional subtitle
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        gap = Cm(0.3)
        quad_w = (sw - int(Cm(3.0)) - int(gap)) // 2
        quad_h = (int(LGDimensions.SLIDE_HEIGHT) - y_pos - int(Cm(0.5)) - int(gap)) // 2
        left_x = Cm(1.5)
        right_x = int(left_x) + quad_w + int(gap)
        top_y = y_pos
        bottom_y = top_y + quad_h + int(gap)

        quadrants = [
            ("Strengths (강점)", strengths or [], left_x, top_y, self.colors.RED),
            ("Weaknesses (약점)", weaknesses or [], right_x, top_y, self.colors.CHARCOAL),
            ("Opportunities (기회)", opportunities or [], left_x, bottom_y, self.colors.GREEN),
            ("Threats (위협)", threats or [], right_x, bottom_y, self.colors.ORANGE),
        ]

        for q_title, items, qx, qy, color in quadrants:
            # Header bar (rounded)
            self.add_box(
                slide, qx, qy, quad_w, Cm(1.3),
                text=q_title, bg_color=color, text_color=self.colors.WHITE,
                font_size=LGTypography.BODY, bold=True, rounded=True
            )

            # Content area (WHITE rounded)
            content_box = self.add_box(
                slide, qx, qy + int(Cm(1.3)), quad_w, quad_h - int(Cm(1.3)),
                bg_color=self.colors.WHITE, rounded=True
            )

            # Items
            if items:
                items_box = slide.shapes.add_textbox(
                    qx + Cm(0.5), qy + int(Cm(1.8)),
                    quad_w - Cm(1.0), quad_h - int(Cm(2.3))
                )
                tf = items_box.text_frame
                tf.word_wrap = True
                items_box.fill.background()
                items_box.line.fill.background()

                for j, item in enumerate(items):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.space_after = Pt(3)
                    pPr = p._p.get_or_add_pPr()
                    buChar = _make_oxml_element('a:buChar')
                    buChar.set('char', '\u2022')
                    pPr.append(buChar)
                    run = p.add_run()
                    run.text = item
                    self._set_font(run, size=LGTypography.BODY_SMALL,
                                  color=self.colors.BLACK)

        return slide

    # ─────────────────────────────────────────
    # Utility Methods
    # ─────────────────────────────────────────

    def add_box(self, slide, left, top, width, height, text="",
                bg_color=None, text_color=None, font_size=None,
                bold=False, alignment=None, border_color=None,
                shadow=True, rounded=True, shadow_type="card"):
        """
        Add a styled content box to a slide with template-matching shadow.

        Args:
            shadow: Whether to add shadow (default True).
            rounded: Whether to use rounded rectangle (default True).
            shadow_type: "card" (ambient), "icon", "content" (directional).

        Returns the shape object.
        """
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        box = slide.shapes.add_shape(
            shape_type,
            left, top, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color or self.colors.WHITE

        # Border: explicit color if given, otherwise transparent (shadow provides depth)
        if border_color:
            box.line.color.rgb = border_color
            box.line.width = Pt(1)
        else:
            box.line.fill.background()

        # Set corner radius for rounded rects (template uses ~4526-14523)
        if rounded:
            spPr = box._element.spPr
            prstGeom = spPr.find(qn('a:prstGeom'))
            if prstGeom is not None:
                avLst = prstGeom.find(qn('a:avLst'))
                if avLst is None:
                    avLst = _make_oxml_element('a:avLst')
                    prstGeom.append(avLst)
                # Clear existing
                for gd in avLst.findall(qn('a:gd')):
                    avLst.remove(gd)
                gd = _make_oxml_element('a:gd')
                gd.set('name', 'adj')
                gd.set('fmla', 'val 8000')
                avLst.append(gd)

        if text:
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Cm(0.15)
            tf.margin_right = Cm(0.15)
            tf.margin_top = Cm(0.05)
            tf.margin_bottom = Cm(0.05)

            p = tf.paragraphs[0]
            p.alignment = alignment or PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            self._set_font(run,
                          size=font_size or LGTypography.BODY,
                          bold=bold,
                          color=text_color or self.colors.BLACK)

        # Add shadow for visual depth (matching L-Style template)
        if shadow:
            self._add_shadow(box, shadow_type=shadow_type)

        return box

    def add_label_badge(self, slide, left, top, text, width=None,
                        bg_color=None, text_color=None):
        """
        Add a small label badge (category tag) using rounded rectangle.
        e.g., DARK_RED background with white text for section labels.
        """
        w = width or Cm(4)
        badge = self.add_box(
            slide, left, top, w, Cm(0.8),
            text=text,
            bg_color=bg_color or self.colors.DARK_RED,
            text_color=text_color or self.colors.WHITE,
            font_size=Pt(9),
            bold=True,
            alignment=PP_ALIGN.CENTER,
            rounded=True
        )
        return badge

    # ─────────────────────────────────────────
    # Advanced Templates
    # ─────────────────────────────────────────

    def add_architecture(self, title, section="", subtitle="",
                         columns=None, rows=None, chapter=""):
        """
        Add a multi-column architecture/system diagram slide — complex grid
        layout with header row, row groups, and area descriptions.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Description text (can be multi-line)
            columns: List of column header strings
            rows: List of row dicts:
                [
                    {
                        "label": "Process User Input",
                        "cells": ["셀1", "셀2", ...],
                        "highlight": False
                    },
                    ...
                ]
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle="", chapter=chapter)
        y_pos = int(y_pos)

        if subtitle:
            # Support multi-line subtitle below chrome
            sub_box = slide.shapes.add_textbox(
                Cm(1.5), y_pos, Cm(31), Cm(1.5)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            sub_box.fill.background()
            sub_box.line.fill.background()
            for i, line in enumerate(subtitle.split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = line
                self._set_font(run, size=LGTypography.SUBTITLE,
                              color=self.colors.DARK_GRAY, weight="semibold")
            y_pos += int(Cm(1.8))

        columns = columns or []
        rows = rows or []
        if not columns:
            return slide

        num_cols = len(columns)
        grid_left = Cm(1.5)
        grid_width = sw - int(Cm(2.0))

        # Column widths: first column slightly narrower for labels
        first_col_w = int(grid_width * 0.18)
        remaining_w = grid_width - first_col_w
        col_w = remaining_w // (num_cols - 1) if num_cols > 1 else remaining_w

        # Header row
        header_h = Cm(1.0)
        x = int(grid_left)
        for ci, col_name in enumerate(columns):
            w = first_col_w if ci == 0 else col_w
            self.add_box(slide, x, y_pos, w, int(header_h),
                        text=col_name,
                        bg_color=self.colors.CHARCOAL,
                        text_color=self.colors.WHITE,
                        font_size=Pt(10), bold=True, shadow=False,
                        rounded=False)
            x += w
        y_pos += int(Cm(1.1))

        # Content rows
        for row in rows:
            label = row.get("label", "")
            cells = row.get("cells", [])
            highlight = row.get("highlight", False)

            # Ensure cells count matches data columns (pad/truncate)
            expected_cells = num_cols - 1
            cells = (cells + [''] * expected_cells)[:expected_cells]

            # Calculate row height based on content
            max_lines = 1
            for cell_text in cells:
                if cell_text:
                    lines = cell_text.count('\n') + 1
                    max_lines = max(max_lines, lines)
            row_h = max(int(Cm(1.2)), int(Cm(0.5 + max_lines * 0.5)))

            x = int(grid_left)
            # First column: label
            label_bg = self.colors.RED if highlight else self.colors.LIGHT_GRAY
            label_tc = self.colors.WHITE if highlight else self.colors.BLACK
            self.add_box(slide, x, y_pos, first_col_w, row_h,
                        text=label,
                        bg_color=label_bg, text_color=label_tc,
                        font_size=Pt(9), bold=True, shadow=False,
                        rounded=False)
            x += first_col_w

            # Remaining columns
            for ci, cell_text in enumerate(cells):
                w = col_w
                if cell_text:
                    # Create box with multi-line text
                    box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        x, y_pos, w, row_h
                    )
                    box.fill.solid()
                    box.fill.fore_color.rgb = self.colors.WHITE
                    box.line.color.rgb = self.colors.BORDER_GRAY
                    box.line.width = Pt(0.5)

                    tf = box.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                    for li, line in enumerate(cell_text.split('\n')):
                        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                        p.alignment = PP_ALIGN.CENTER
                        run = p.add_run()
                        run.text = line
                        # First line bold if multiple lines (acts as sub-header)
                        is_bold = (li == 0 and cell_text.count('\n') > 0)
                        self._set_font(run, size=Pt(9),
                                      bold=is_bold,
                                      color=self.colors.BLACK)

                else:
                    # Empty cell - light border only
                    box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        x, y_pos, w, row_h
                    )
                    box.fill.background()
                    box.line.color.rgb = self.colors.BORDER_GRAY
                    box.line.width = Pt(0.5)
                x += w

            y_pos += row_h + int(Cm(0.1))

        return slide

    def add_strategy_pillars(self, title, section="", subtitle="",
                             pillars=None, chapter=""):
        """
        Add a strategy pillars slide — 3 to 5 vertical columns showing
        strategic focus areas.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Description text
            pillars: List of dicts:
                [
                    {
                        "header": "디지털 전환",
                        "header_color": "#A50034",
                        "items": ["AI/ML 기반 공정 최적화", "디지털 트윈 구축", ...]
                    },
                    ...
                ]
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        pillars = pillars or []
        if not pillars:
            return slide

        num = len(pillars)
        gap = Cm(0.3)
        total_w = sw - int(Cm(2.0))
        pillar_w = (total_w - int(gap) * (num - 1)) // num
        header_h = Cm(1.2)
        body_h = int(LGDimensions.SLIDE_HEIGHT) - y_pos - int(header_h) - int(Cm(0.5))

        x = Cm(1.5)
        for pillar in pillars:
            header = pillar.get("header", "") or pillar.get("title", "")
            items = pillar.get("items", [])
            hdr_color_hex = pillar.get("header_color", None)

            if hdr_color_hex:
                hdr_color = RGBColor(
                    int(hdr_color_hex[1:3], 16),
                    int(hdr_color_hex[3:5], 16),
                    int(hdr_color_hex[5:7], 16)
                )
            else:
                hdr_color = self.colors.CHARCOAL

            # Header (rounded)
            self.add_box(slide, int(x), y_pos, pillar_w, int(header_h),
                        text=header, bg_color=hdr_color,
                        text_color=self.colors.WHITE,
                        font_size=Pt(14), bold=True, shadow=False,
                        rounded=True)

            # Body (WHITE rounded)
            body_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(x), y_pos + int(header_h), pillar_w, body_h
            )
            body_box.fill.solid()
            body_box.fill.fore_color.rgb = self.colors.WHITE
            body_box.line.color.rgb = self.colors.BORDER_GRAY
            body_box.line.width = Pt(0.5)

            tf = body_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Cm(0.3)
            tf.margin_right = Cm(0.3)
            tf.margin_top = Cm(0.3)

            for j, item in enumerate(items):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(3)
                # Bullet
                pPr = p._p.get_or_add_pPr()
                buChar = _make_oxml_element('a:buChar')
                buChar.set('char', '\u2022')
                pPr.append(buChar)
                pPr.set('marL', str(int(Cm(0.5))))
                pPr.set('indent', str(int(Cm(-0.4))))

                run = p.add_run()
                run.text = item
                self._set_font(run, size=LGTypography.BODY,
                              color=self.colors.BLACK)

            self._add_shadow(body_box)
            x = int(x) + pillar_w + int(gap)

        return slide

    def add_risk_matrix(self, title, section="", subtitle="",
                        risks=None, x_label="영향도 (Impact)",
                        y_label="발생 가능성 (Likelihood)", chapter=""):
        """
        Add a risk assessment matrix slide — 3x3 grid with color-coded risk
        levels.

        Args:
            title: Slide title
            section: Section name
            subtitle: Description text
            risks: List of risk items to place on the matrix:
                [
                    {"name": "데이터 유출", "likelihood": 2, "impact": 3},
                    {"name": "인력 이탈", "likelihood": 1, "impact": 2},
                ]
                where likelihood and impact are 1 (Low), 2 (Medium), 3 (High)
            x_label: Label for the X axis (impact)
            y_label: Label for the Y axis (likelihood)
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        risks = risks or []

        # Grid layout
        grid_left = Cm(4.5)
        grid_top = y_pos + int(Cm(0.5))
        cell_size = Cm(4.0)
        grid_gap = Cm(0.15)

        # Color coding: [row][col] → (likelihood high→low, impact low→high)
        colors_3x3 = [
            # impact Low, Med, High  (likelihood High = top row)
            [RGBColor(0xD4, 0x76, 0x0A), RGBColor(0xA5, 0x00, 0x34), RGBColor(0xA5, 0x00, 0x34)],
            [RGBColor(0x2E, 0x7D, 0x32), RGBColor(0xD4, 0x76, 0x0A), RGBColor(0xA5, 0x00, 0x34)],
            [RGBColor(0x2E, 0x7D, 0x32), RGBColor(0x2E, 0x7D, 0x32), RGBColor(0xD4, 0x76, 0x0A)],
        ]
        labels_3x3 = [
            ["중간", "높음", "매우 높음"],
            ["낮음", "중간", "높음"],
            ["매우 낮음", "낮음", "중간"],
        ]

        # Y axis label
        self._add_textbox(
            slide, Cm(1.5), grid_top, Cm(2.5), int(cell_size) * 3 + int(grid_gap) * 2,
            text=y_label, size=Pt(10), bold=True, color=self.colors.DARK_GRAY
        )

        # Y axis level labels
        y_levels = ["High", "Medium", "Low"]
        for r in range(3):
            cy = grid_top + r * (int(cell_size) + int(grid_gap))
            self._add_textbox(
                slide, Cm(2.8), cy, Cm(1.5), int(cell_size),
                text=y_levels[r], size=Pt(9), color=self.colors.MEDIUM_GRAY,
                vertical=MSO_ANCHOR.MIDDLE
            )

        # X axis label
        self._add_textbox(
            slide, int(grid_left), grid_top + 3 * (int(cell_size) + int(grid_gap)) + int(Cm(0.3)),
            int(cell_size) * 3 + int(grid_gap) * 2, Cm(1.0),
            text=x_label, size=Pt(10), bold=True, color=self.colors.DARK_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # X axis level labels
        x_levels = ["Low", "Medium", "High"]
        for c in range(3):
            cx = int(grid_left) + c * (int(cell_size) + int(grid_gap))
            self._add_textbox(
                slide, cx,
                grid_top + 3 * (int(cell_size) + int(grid_gap)),
                int(cell_size), Cm(0.8),
                text=x_levels[c], size=Pt(9), color=self.colors.MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        # Draw grid cells
        for r in range(3):
            for c in range(3):
                cx = int(grid_left) + c * (int(cell_size) + int(grid_gap))
                cy = grid_top + r * (int(cell_size) + int(grid_gap))
                cell_color = colors_3x3[r][c]

                box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, int(cell_size), int(cell_size)
                )
                box.fill.solid()
                box.fill.fore_color.rgb = cell_color
                box.line.fill.background()

                # Muted label inside
                tf = box.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.BOTTOM
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT
                run = p.add_run()
                run.text = labels_3x3[r][c]
                self._set_font(run, size=Pt(8), color=self.colors.WHITE)

        # Place risk items
        for risk in risks:
            name = risk.get("name", "")
            lk = risk.get("likelihood", 1)  # 1-3
            im = risk.get("impact", 1)       # 1-3
            # Map to grid: likelihood 3=top(row0), impact 1=left(col0)
            r = 3 - lk
            c = im - 1
            cx = int(grid_left) + c * (int(cell_size) + int(grid_gap)) + int(Cm(0.3))
            cy = grid_top + r * (int(cell_size) + int(grid_gap)) + int(Cm(0.3))

            # Risk item tag
            self.add_box(
                slide, cx, cy,
                int(cell_size) - int(Cm(0.6)), Cm(1.0),
                text=name, bg_color=self.colors.WHITE,
                text_color=self.colors.BLACK,
                font_size=Pt(8), bold=True,
                shadow=True, rounded=True
            )

        # Legend on the right side
        legend_x = int(grid_left) + 3 * (int(cell_size) + int(grid_gap)) + int(Cm(1.5))
        legend_y = grid_top
        self._add_textbox(
            slide, legend_x, legend_y, Cm(5), Cm(0.8),
            text="위험 수준", size=Pt(10), bold=True, color=self.colors.BLACK
        )
        legend_items = [
            (self.colors.RED, "높음 (High Risk)"),
            (RGBColor(0xD4, 0x76, 0x0A), "중간 (Medium Risk)"),
            (self.colors.GREEN, "낮음 (Low Risk)"),
        ]
        for li, (color, label) in enumerate(legend_items):
            ly = legend_y + int(Cm(1.0)) + li * int(Cm(1.0))
            # Color swatch
            swatch = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, legend_x, ly, Cm(0.6), Cm(0.6)
            )
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = color
            swatch.line.fill.background()
            # Label
            self._add_textbox(
                slide, legend_x + int(Cm(1.0)), ly, Cm(5), Cm(0.6),
                text=label, size=Pt(9), color=self.colors.BLACK
            )

        return slide

    def add_financial_summary(self, title, section="", subtitle="",
                              categories=None, total_label="합계",
                              chapter=""):
        """
        Add a financial/budget summary slide — structured table with
        category groups, subtotals and a grand total row.

        Args:
            title: Slide title
            section: Section name
            subtitle: Description text
            categories: List of category dicts (see original docstring)
            total_label: Label for grand total row
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        categories = categories or []
        if not categories:
            return slide

        # Determine year columns from first category
        sample_item = categories[0]["items"][0] if categories[0].get("items") else {}
        year_keys = [k for k in sample_item.keys() if k != "name"]

        # Build table data
        headers = ["구분", "항목"] + [f"{y}년" if len(y) == 4 else y for y in year_keys]
        num_cols = len(headers)

        # Count total rows
        total_rows = 1  # header
        for cat in categories:
            total_rows += len(cat.get("items", []))
            if cat.get("subtotal"):
                total_rows += 1  # subtotal row
        total_rows += 1  # grand total

        table_left = Cm(1.5)
        table_width = sw - int(Cm(2.0))
        row_h = Cm(0.9)

        # Calculate column widths first, then set table width to exact sum
        cat_col_w = int(table_width * 0.12)
        item_col_w = int(table_width * 0.22)
        year_col_w = (table_width - cat_col_w - item_col_w) // max(len(year_keys), 1)
        exact_width = cat_col_w + item_col_w + year_col_w * len(year_keys)

        table_shape = slide.shapes.add_table(
            total_rows, num_cols,
            int(table_left), y_pos,
            exact_width, int(row_h) * total_rows
        )
        table = table_shape.table

        # Clear default table style to prevent rendering conflicts
        tbl = table._tbl
        tblPr = tbl.tblPr
        if tblPr is not None:
            tblPr.set('bandRow', '0')
            tblPr.set('bandCol', '0')
            tblPr.set('firstRow', '0')
            tblPr.set('lastRow', '0')
            tblPr.set('firstCol', '0')
            tblPr.set('lastCol', '0')

        # Column widths
        table.columns[0].width = cat_col_w
        table.columns[1].width = item_col_w
        for ci in range(2, num_cols):
            table.columns[ci].width = year_col_w

        # Header row
        for ci, header in enumerate(headers):
            self._format_table_cell(table.cell(0, ci), text=header,
                                   is_header=True)
            self._set_cell_border(table.cell(0, ci))

        # Data rows
        row_idx = 1
        for cat in categories:
            cat_name = cat.get("name", "")
            items = cat.get("items", [])
            subtotal = cat.get("subtotal", None)

            # Category items
            for ii, item in enumerate(items):
                # Category column (merge visual: only first row shows name)
                if ii == 0:
                    self._format_table_cell(
                        table.cell(row_idx, 0), text=cat_name,
                        bold=True, bg_color=self.colors.LIGHT_GRAY,
                        alignment=PP_ALIGN.CENTER
                    )
                else:
                    self._format_table_cell(
                        table.cell(row_idx, 0), text="",
                        bg_color=self.colors.LIGHT_GRAY
                    )

                # Item name
                self._format_table_cell(
                    table.cell(row_idx, 1), text=item.get("name", ""),
                    alignment=PP_ALIGN.LEFT
                )

                # Year values
                for yi, yk in enumerate(year_keys):
                    self._format_table_cell(
                        table.cell(row_idx, 2 + yi),
                        text=str(item.get(yk, "")),
                        alignment=PP_ALIGN.RIGHT
                    )

                for ci in range(num_cols):
                    self._set_cell_border(table.cell(row_idx, ci))
                row_idx += 1

            # Subtotal row
            if subtotal:
                self._format_table_cell(
                    table.cell(row_idx, 0), text="",
                    bg_color=RGBColor(0xE8, 0xE8, 0xE8)
                )
                self._format_table_cell(
                    table.cell(row_idx, 1), text="소계",
                    bold=True, alignment=PP_ALIGN.CENTER,
                    bg_color=RGBColor(0xE8, 0xE8, 0xE8)
                )
                for yi, yk in enumerate(year_keys):
                    self._format_table_cell(
                        table.cell(row_idx, 2 + yi),
                        text=str(subtotal.get(yk, "")),
                        bold=True, alignment=PP_ALIGN.RIGHT,
                        bg_color=RGBColor(0xE8, 0xE8, 0xE8)
                    )
                for ci in range(num_cols):
                    self._set_cell_border(table.cell(row_idx, ci))
                row_idx += 1

        # Grand total row
        self._format_table_cell(
            table.cell(row_idx, 0), text="",
            is_header=True, bg_color=self.colors.CHARCOAL
        )
        self._format_table_cell(
            table.cell(row_idx, 1), text=total_label,
            is_header=True, bg_color=self.colors.CHARCOAL,
            alignment=PP_ALIGN.CENTER
        )
        # Sum subtotals for grand total
        for yi, yk in enumerate(year_keys):
            grand = 0
            for cat in categories:
                st = cat.get("subtotal", {})
                try:
                    grand += float(st.get(yk, 0))
                except (ValueError, TypeError):
                    pass
            self._format_table_cell(
                table.cell(row_idx, 2 + yi),
                text=f"{grand:.1f}",
                is_header=True, bg_color=self.colors.CHARCOAL,
                alignment=PP_ALIGN.RIGHT
            )
        for ci in range(num_cols):
            self._set_cell_border(table.cell(row_idx, ci), color="3C3C3C")

        # Merge category cells vertically
        for cat in categories:
            pass  # Visual merging done by leaving subsequent rows empty

        return slide

    def add_milestone_tracker(self, title, section="", subtitle="",
                              phases=None, chapter=""):
        """
        Add a milestone/progress tracker slide — horizontal chevron phases
        with detailed milestones underneath.

        Args:
            title: Slide title
            section: Section name
            subtitle: Description text
            phases: List of phase dicts (see original docstring)
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        phases = phases or []
        if not phases:
            return slide

        num = len(phases)
        gap = Cm(0.3)
        total_w = sw - int(Cm(2.0))
        phase_w = (total_w - int(gap) * (num - 1)) // num
        chevron_h = Cm(1.5)
        status_colors = {
            "completed": self.colors.CHARCOAL,
            "in_progress": self.colors.RED,
            "planned": self.colors.MEDIUM_GRAY,
        }

        # Phase chevrons
        x = Cm(1.5)
        for pi, phase in enumerate(phases):
            status = phase.get("status", "planned")
            bg = status_colors.get(status, self.colors.MEDIUM_GRAY)

            # Chevron header (rounded)
            chevron = self.add_box(
                slide, int(x), y_pos, phase_w, int(chevron_h),
                text=phase.get("name", ""),
                bg_color=bg, text_color=self.colors.WHITE,
                font_size=Pt(10), bold=True, shadow=False,
                rounded=True
            )

            # Period label below chevron
            period_y = y_pos + int(chevron_h) + int(Cm(0.1))
            self._add_textbox(
                slide, int(x), period_y, phase_w, Cm(0.6),
                text=phase.get("period", ""),
                size=Pt(9), color=bg,
                alignment=PP_ALIGN.CENTER
            )

            # Status badge
            badge_y = period_y + int(Cm(0.7))
            status_text = {"completed": "완료", "in_progress": "진행중",
                          "planned": "예정"}.get(status, "")
            status_badge_bg = {
                "completed": self.colors.GREEN,
                "in_progress": self.colors.RED,
                "planned": self.colors.BORDER_GRAY,
            }.get(status, self.colors.BORDER_GRAY)
            self.add_box(
                slide, int(x) + phase_w // 4, badge_y,
                phase_w // 2, Cm(0.5),
                text=status_text, bg_color=status_badge_bg,
                text_color=self.colors.WHITE,
                font_size=Pt(8), bold=True, shadow=False,
                rounded=True
            )

            # Milestones below (WHITE rounded box)
            milestones_y = badge_y + int(Cm(0.8))
            milestones = phase.get("milestones", [])
            if milestones:
                ms_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    int(x), milestones_y, phase_w,
                    int(LGDimensions.SLIDE_HEIGHT) - milestones_y - int(Cm(1.0))
                )
                tf = ms_box.text_frame
                tf.word_wrap = True
                ms_box.fill.solid()
                ms_box.fill.fore_color.rgb = self.colors.WHITE
                ms_box.line.color.rgb = self.colors.BORDER_GRAY
                ms_box.line.width = Pt(0.5)
                tf.margin_left = Cm(0.3)
                tf.margin_right = Cm(0.2)
                tf.margin_top = Cm(0.3)

                for j, ms in enumerate(milestones):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.space_after = Pt(3)
                    pPr = p._p.get_or_add_pPr()
                    buChar = _make_oxml_element('a:buChar')
                    buChar.set('char', '\u2713' if status == "completed" else '\u25b8')
                    pPr.append(buChar)
                    pPr.set('marL', str(int(Cm(0.5))))
                    pPr.set('indent', str(int(Cm(-0.4))))
                    run = p.add_run()
                    run.text = ms
                    color = self.colors.BLACK if status != "planned" else self.colors.MEDIUM_GRAY
                    self._set_font(run, size=Pt(10), color=color)

                self._add_shadow(ms_box)

            x = int(x) + phase_w + int(gap)

        return slide

    def add_comparison_cards(self, title, section="", subtitle="",
                             cards=None, chapter=""):
        """
        Add a comparison cards slide — side-by-side cards for comparing
        options, vendors, solutions, or before/after states.

        Args:
            title: Slide title
            section: Section name
            subtitle: Description text
            cards: List of card dicts (2-4 cards, see original docstring)
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        cards = cards or []
        if not cards:
            return slide

        num = len(cards)
        gap = Cm(0.3)
        total_w = sw - int(Cm(2.0))
        card_w = (total_w - int(gap) * (num - 1)) // num
        card_h = int(LGDimensions.SLIDE_HEIGHT) - y_pos - int(Cm(0.5))
        header_h = Cm(1.0)

        x = Cm(1.5)
        for card in cards:
            header = card.get("header", "")
            highlight = card.get("highlight", False)
            metrics = card.get("metrics", [])
            pros = card.get("pros", [])
            cons = card.get("cons", [])

            hdr_color_hex = card.get("header_color", None)
            if hdr_color_hex:
                hdr_bg = RGBColor(
                    int(hdr_color_hex[1:3], 16),
                    int(hdr_color_hex[3:5], 16),
                    int(hdr_color_hex[5:7], 16)
                )
            else:
                hdr_bg = self.colors.RED if highlight else self.colors.CHARCOAL

            # Header with optional badge (rounded)
            hdr_text = f"★ {header} (추천)" if highlight else header
            self.add_box(
                slide, int(x), y_pos, card_w, int(header_h),
                text=hdr_text, bg_color=hdr_bg,
                text_color=self.colors.WHITE,
                font_size=Pt(11), bold=True, shadow=False,
                rounded=True
            )

            # Card body (WHITE rounded)
            body_top = y_pos + int(header_h)
            body_h = card_h - int(header_h)
            border = self.colors.RED if highlight else self.colors.BORDER_GRAY

            body = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(x), body_top, card_w, body_h
            )
            body.fill.solid()
            body.fill.fore_color.rgb = self.colors.WHITE
            body.line.color.rgb = border
            body.line.width = Pt(1.5 if highlight else 0.75)

            tf = body.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Cm(0.4)
            tf.margin_right = Cm(0.4)
            tf.margin_top = Cm(0.3)

            # Metrics
            for mi, metric in enumerate(metrics):
                p = tf.paragraphs[0] if mi == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = f"{metric['label']}: "
                self._set_font(run, size=Pt(10), color=self.colors.MEDIUM_GRAY)
                run2 = p.add_run()
                run2.text = metric['value']
                self._set_font(run2, size=Pt(10), bold=True, color=self.colors.BLACK)

            # Separator
            if metrics and (pros or cons):
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(4)
                p.space_after = Pt(4)
                run = p.add_run()
                run.text = "\u2500" * 20
                self._set_font(run, size=Pt(6), color=self.colors.BORDER_GRAY)

            # Pros
            if pros:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = "장점"
                self._set_font(run, size=Pt(9), bold=True, color=self.colors.GREEN)
                for item in pros:
                    p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(2)
                    run = p.add_run()
                    run.text = f"  + {item}"
                    self._set_font(run, size=Pt(9), color=self.colors.BLACK)

            # Cons
            if cons:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(4)
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = "단점"
                self._set_font(run, size=Pt(9), bold=True,
                              color=RGBColor(0xA5, 0x00, 0x34))
                for item in cons:
                    p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(2)
                    run = p.add_run()
                    run.text = f"  - {item}"
                    self._set_font(run, size=Pt(9), color=self.colors.BLACK)

            self._add_shadow(body)
            x = int(x) + card_w + int(gap)

        return slide

    # ─────────────────────────────────────────
    # New Slide Types
    # ─────────────────────────────────────────

    def add_gantt_chart(self, title, section="", subtitle="", tasks=None,
                        start_date="", months=6, chapter=""):
        """
        Add a Gantt chart slide showing project task timelines.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            tasks: List of task dicts (see original docstring)
            start_date: Label for the start period (e.g., "2025.01")
            months: Total number of months to display (default 6)
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        tasks = tasks or []
        if not tasks:
            return slide

        # Layout dimensions
        name_col_width = Cm(6)
        chart_left = int(Cm(1.5)) + int(name_col_width)
        chart_width = sw - chart_left - int(Cm(0.5))
        month_width = chart_width // months
        row_height = Cm(1.0)
        bar_height = Cm(0.6)
        header_height = Cm(0.8)

        # Month header row
        for m in range(months):
            mx = chart_left + month_width * m
            # Header cell
            self.add_box(
                slide, mx, y_pos, month_width, int(header_height),
                text=f"M{m + 1}" if not start_date else "",
                bg_color=self.colors.CHARCOAL, text_color=self.colors.WHITE,
                font_size=Pt(9), bold=True, shadow=False, rounded=False
            )
            # Vertical grid line
            if m > 0:
                grid_line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    mx, y_pos + int(header_height),
                    Pt(0.5), int(row_height) * len(tasks)
                )
                grid_line.fill.solid()
                grid_line.fill.fore_color.rgb = self.colors.BORDER_GRAY
                grid_line.line.fill.background()

        # Month labels if start_date provided
        if start_date:
            try:
                base_year = int(start_date.split('.')[0])
                base_month = int(start_date.split('.')[1])
            except (ValueError, IndexError):
                base_year, base_month = 2025, 1
            for m in range(months):
                cm = ((base_month - 1 + m) % 12) + 1
                cy = base_year + (base_month - 1 + m) // 12
                mx = chart_left + month_width * m
                self._add_textbox(
                    slide, mx, y_pos,
                    month_width, int(header_height),
                    text=f"{cy}.{cm:02d}",
                    size=Pt(8), bold=True, color=self.colors.WHITE,
                    alignment=PP_ALIGN.CENTER
                )

        # Task name header
        self.add_box(
            slide, Cm(1.5), y_pos, int(name_col_width), int(header_height),
            text="작업 항목", bg_color=self.colors.CHARCOAL,
            text_color=self.colors.WHITE, font_size=Pt(9), bold=True,
            shadow=False, rounded=False
        )

        y_pos = y_pos + int(header_height)

        # Task rows
        for ti, task in enumerate(tasks):
            task_name = task.get("name", "")
            t_start = task.get("start", 0)
            t_dur = task.get("duration", 1)
            t_progress = task.get("progress", 0)
            t_color_hex = task.get("color", "#A50034")

            try:
                tr, tg, tb = int(t_color_hex[1:3], 16), int(t_color_hex[3:5], 16), int(t_color_hex[5:7], 16)
                bar_color = RGBColor(tr, tg, tb)
            except (ValueError, IndexError):
                bar_color = self.colors.RED

            row_y = y_pos + int(row_height) * ti

            # Alternating row background
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                chart_left, row_y, chart_width, int(row_height)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = self.colors.WHITE if ti % 2 == 0 else self.colors.LIGHT_GRAY
            row_bg.line.color.rgb = self.colors.RED
            row_bg.line.width = Pt(1)
            # Set corner radius (template: val 8000)
            _spPr = row_bg._element.spPr
            _pg = _spPr.find(qn('a:prstGeom'))
            if _pg is not None:
                _al = _pg.find(qn('a:avLst'))
                if _al is None:
                    _al = _make_oxml_element('a:avLst')
                    _pg.append(_al)
                for _gd in _al.findall(qn('a:gd')):
                    _al.remove(_gd)
                _gd = _make_oxml_element('a:gd')
                _gd.set('name', 'adj')
                _gd.set('fmla', 'val 8000')
                _al.append(_gd)

            # Task name
            self._add_textbox(
                slide, Cm(1.5), row_y,
                int(name_col_width), int(row_height),
                text=task_name, size=Pt(10), color=self.colors.BLACK,
                alignment=PP_ALIGN.LEFT,
                vertical=MSO_ANCHOR.MIDDLE
            )

            # Bar background (total duration)
            bar_x = chart_left + month_width * t_start
            bar_w = month_width * t_dur
            bar_y = row_y + (int(row_height) - int(bar_height)) // 2

            bg_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_x, bar_y, bar_w, int(bar_height)
            )
            bg_bar.fill.solid()
            bg_bar.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            bg_bar.line.fill.background()

            # Progress fill
            if t_progress > 0:
                fill_w = max(int(bar_w * t_progress / 100), int(Cm(0.1)))
                fill_bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    bar_x, bar_y, fill_w, int(bar_height)
                )
                fill_bar.fill.solid()
                fill_bar.fill.fore_color.rgb = bar_color
                fill_bar.line.fill.background()

            # Progress label
            if t_progress > 0:
                self._add_textbox(
                    slide, bar_x, bar_y,
                    bar_w, int(bar_height),
                    text=f"{t_progress}%",
                    size=Pt(7), bold=True, color=self.colors.WHITE,
                    alignment=PP_ALIGN.CENTER,
                    vertical=MSO_ANCHOR.MIDDLE
                )

        return slide

    def add_org_chart(self, title, section="", subtitle="", org_data=None,
                      chapter=""):
        """
        Add an organizational chart slide with hierarchical boxes and connectors.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            org_data: Dict representing the org hierarchy (see original docstring)
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        if not org_data:
            return slide

        # Calculate tree structure
        def count_leaves(node):
            children = node.get("children", [])
            if not children:
                return 1
            return sum(count_leaves(c) for c in children)

        def get_depth(node):
            children = node.get("children", [])
            if not children:
                return 1
            return 1 + max(get_depth(c) for c in children)

        depth = get_depth(org_data)
        total_leaves = count_leaves(org_data)

        # Box dimensions
        box_w = Cm(4.5)
        box_h = Cm(1.6)
        h_gap = Cm(0.5)
        v_gap = Cm(1.2)
        level_height = int(box_h) + int(v_gap)

        available_width = sw - int(Cm(2.0))
        content_top = y_pos

        # Recursive drawing with position tracking
        def draw_node(node, level, cx, is_root=False):
            """Draw a node at center_x=cx, returns (center_x, top_y, bottom_y)."""
            ny = content_top + level * level_height
            nx = cx - int(box_w) // 2

            # Box styling (rounded)
            bg = self.colors.CHARCOAL if is_root else self.colors.WHITE
            tc = self.colors.WHITE if is_root else self.colors.BLACK

            box = self.add_box(
                slide, nx, ny, int(box_w), int(box_h),
                bg_color=bg, shadow=True, rounded=True
            )

            # Name (bold)
            name = node.get("name", "")
            node_title = node.get("title", "")
            name_box = slide.shapes.add_textbox(
                nx, ny, int(box_w), int(box_h) // 2
            )
            name_tf = name_box.text_frame
            name_tf.word_wrap = True
            name_box.fill.background()
            name_box.line.fill.background()
            p = name_tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = name
            self._set_font(run, size=Pt(10), bold=True, color=tc)

            # Title
            if node_title:
                title_box = slide.shapes.add_textbox(
                    nx, ny + int(box_h) // 2, int(box_w), int(box_h) // 2
                )
                title_tf = title_box.text_frame
                title_tf.word_wrap = True
                title_box.fill.background()
                title_box.line.fill.background()
                p = title_tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = node_title
                self._set_font(run, size=Pt(8), color=tc if is_root else self.colors.MEDIUM_GRAY)

            children = node.get("children", [])
            if children:
                # Calculate children positions
                child_leaves = [count_leaves(c) for c in children]
                total_child_leaves = sum(child_leaves)

                # Minimum width per leaf
                min_leaf_w = int(box_w) + int(h_gap)
                total_children_width = total_child_leaves * min_leaf_w

                # Center children under parent
                child_start_x = cx - total_children_width // 2
                child_centers = []

                offset = 0
                for ci, child in enumerate(children):
                    cl = child_leaves[ci]
                    child_cx = child_start_x + offset + (cl * min_leaf_w) // 2
                    child_centers.append(child_cx)
                    draw_node(child, level + 1, child_cx)
                    offset += cl * min_leaf_w

                # Draw connectors
                parent_bottom_y = ny + int(box_h)
                child_top_y = content_top + (level + 1) * level_height
                mid_y = (parent_bottom_y + child_top_y) // 2

                # Vertical line from parent down
                v_line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    cx - Pt(1), parent_bottom_y,
                    Pt(2), mid_y - parent_bottom_y
                )
                v_line.fill.solid()
                v_line.fill.fore_color.rgb = self.colors.BORDER_GRAY
                v_line.line.fill.background()

                # Horizontal line connecting children
                if len(child_centers) > 1:
                    h_left = min(child_centers)
                    h_right = max(child_centers)
                    h_line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        h_left, mid_y,
                        h_right - h_left, Pt(2)
                    )
                    h_line.fill.solid()
                    h_line.fill.fore_color.rgb = self.colors.BORDER_GRAY
                    h_line.line.fill.background()

                # Vertical lines from horizontal to each child
                for ccx in child_centers:
                    vc_line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        ccx - Pt(1), mid_y,
                        Pt(2), child_top_y - mid_y
                    )
                    vc_line.fill.solid()
                    vc_line.fill.fore_color.rgb = self.colors.BORDER_GRAY
                    vc_line.line.fill.background()

        # Draw from root centered on slide
        draw_node(org_data, 0, sw // 2, is_root=True)

        return slide

    def add_pyramid(self, title, section="", subtitle="", levels=None,
                    chapter=""):
        """
        Add a pyramid/hierarchy diagram slide with trapezoid layers.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            levels: List of level dicts (top to bottom, see original docstring)
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        levels = levels or []
        if not levels:
            return slide

        num_levels = len(levels)
        available_height = sh - y_pos - int(Cm(0.5))
        level_height = available_height // num_levels
        gap = Cm(0.1)

        # Pyramid center and width range
        center_x = sw // 2 - int(Cm(2))  # Shift left to leave room for descriptions
        min_width = Cm(6)
        max_width = Cm(18)
        desc_left = center_x + int(max_width) // 2 + int(Cm(1.0))

        # Distinct color palette for each level
        _pyramid_palette = [
            "#A50034", "#3C3C3C", "#1565C0", "#2E7D32", "#D4760A", "#7B1FA2"
        ]
        for i, level in enumerate(levels):
            # Calculate width for this level (narrower at top)
            fraction = i / max(num_levels - 1, 1)
            level_w = int(min_width) + int((int(max_width) - int(min_width)) * fraction)
            lx = center_x - level_w // 2
            ly = y_pos + level_height * i

            # Color: use per-level color key, or palette, or fallback
            level_color_hex = level.get("color", None)
            if not level_color_hex:
                level_color_hex = _pyramid_palette[i % len(_pyramid_palette)]
            level_color = RGBColor(
                int(level_color_hex[1:3], 16),
                int(level_color_hex[3:5], 16),
                int(level_color_hex[5:7], 16)
            )

            # Trapezoid shape (rounded rectangle)
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                lx, ly, level_w, level_height - int(gap)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = level_color
            box.line.fill.background()

            # Level label (centered in trapezoid)
            label = level.get("label", "")
            self._add_textbox(
                slide, lx, ly,
                level_w, level_height - int(gap),
                text=label, size=Pt(14), bold=True,
                color=self.colors.WHITE,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

            # Description on the right
            desc = level.get("description", "")
            if desc:
                # Connector line
                conn_y = ly + (level_height - int(gap)) // 2
                conn = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    lx + level_w, conn_y,
                    desc_left - (lx + level_w), Pt(1)
                )
                conn.fill.solid()
                conn.fill.fore_color.rgb = self.colors.BORDER_GRAY
                conn.line.fill.background()

                # Description text
                self._add_textbox(
                    slide, desc_left, ly,
                    sw - desc_left - int(Cm(0.5)), level_height - int(gap),
                    text=desc, size=Pt(10), color=self.colors.DARK_GRAY,
                    alignment=PP_ALIGN.LEFT,
                    vertical=MSO_ANCHOR.MIDDLE
                )

        return slide

    def add_positioning_map(self, title, section="", subtitle="",
                            x_label="", y_label="", items=None,
                            quadrant_labels=None, chapter=""):
        """
        Add a 2D positioning/perceptual map slide with plotted items.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            x_label: Label for X axis
            y_label: Label for Y axis
            items: List of item dicts (see original docstring)
            quadrant_labels: Optional list of 4 labels
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        items = items or []

        # Map area
        map_left = int(Cm(3.5))
        map_top = y_pos + int(Cm(0.5))
        map_size = min(sw - int(Cm(5.0)), sh - map_top - int(Cm(1.5)))
        map_width = map_size
        map_height = map_size

        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            map_left, map_top, map_width, map_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors.LIGHT_GRAY
        bg.line.color.rgb = self.colors.BORDER_GRAY
        bg.line.width = Pt(1)

        # Cross-hairs (axes through center)
        mid_x = map_left + map_width // 2
        mid_y = map_top + map_height // 2

        # Vertical axis
        v_axis = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            mid_x - Pt(1), map_top, Pt(2), map_height
        )
        v_axis.fill.solid()
        v_axis.fill.fore_color.rgb = self.colors.BORDER_GRAY
        v_axis.line.fill.background()

        # Horizontal axis
        h_axis = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            map_left, mid_y - Pt(1), map_width, Pt(2)
        )
        h_axis.fill.solid()
        h_axis.fill.fore_color.rgb = self.colors.BORDER_GRAY
        h_axis.line.fill.background()

        # Axis labels
        if x_label:
            self._add_textbox(
                slide, map_left, map_top + map_height + int(Cm(0.2)),
                map_width, Cm(0.8),
                text=x_label, size=Pt(10), bold=True,
                color=self.colors.DARK_GRAY, alignment=PP_ALIGN.CENTER
            )
        if y_label:
            self._add_textbox(
                slide, map_left - int(Cm(3.0)), map_top,
                Cm(2.5), map_height,
                text=y_label, size=Pt(10), bold=True,
                color=self.colors.DARK_GRAY, alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

        # Quadrant labels
        if quadrant_labels and len(quadrant_labels) >= 4:
            ql_positions = [
                (map_left + int(Cm(0.3)), map_top + int(Cm(0.2))),
                (mid_x + int(Cm(0.3)), map_top + int(Cm(0.2))),
                (map_left + int(Cm(0.3)), mid_y + int(Cm(0.2))),
                (mid_x + int(Cm(0.3)), mid_y + int(Cm(0.2))),
            ]
            for qi, (qx, qy) in enumerate(ql_positions):
                if qi < len(quadrant_labels):
                    self._add_textbox(
                        slide, qx, qy,
                        map_width // 2 - int(Cm(0.6)), Cm(0.6),
                        text=quadrant_labels[qi],
                        size=Pt(8), color=self.colors.MEDIUM_GRAY,
                        alignment=PP_ALIGN.LEFT
                    )

        # Plot items
        size_map = {"small": Cm(1.2), "medium": Cm(1.8), "large": Cm(2.5)}
        for item in items:
            ix = item.get("x", 0.5)
            iy = item.get("y", 0.5)
            i_size = size_map.get(item.get("size", "medium"), Cm(1.8))
            i_name = item.get("name", "")
            i_color_hex = item.get("color", "#A50034")

            try:
                ir, ig, ib = int(i_color_hex[1:3], 16), int(i_color_hex[3:5], 16), int(i_color_hex[5:7], 16)
                i_color = RGBColor(ir, ig, ib)
            except (ValueError, IndexError):
                i_color = self.colors.RED

            # Convert normalized coords to map coords
            px = map_left + int(map_width * ix) - int(i_size) // 2
            py = map_top + int(map_height * (1.0 - iy)) - int(i_size) // 2

            # Bubble
            bubble = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                px, py, int(i_size), int(i_size)
            )
            bubble.fill.solid()
            bubble.fill.fore_color.rgb = i_color

            # Set transparency (40%)
            spPr = bubble._element.spPr
            solidFill = spPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    alpha = _make_oxml_element('a:alpha')
                    alpha.set('val', '60000')  # 60% opaque
                    srgbClr.append(alpha)

            bubble.line.fill.background()

            # Label
            self._add_textbox(
                slide, px, py,
                int(i_size), int(i_size),
                text=i_name, size=Pt(8), bold=True,
                color=self.colors.WHITE,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

        return slide

    def add_keyword_highlight(self, title, section="", subtitle="",
                              keywords=None, description="", chapter=""):
        """
        Add a keyword emphasis/tag-cloud style slide.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            keywords: List of keyword dicts (see original docstring)
            description: Supporting description text below keywords
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        keywords = keywords or []
        if not keywords:
            return slide

        # Keyword cloud area
        cloud_left = Cm(2.0)
        cloud_top = y_pos + int(Cm(0.5))
        cloud_width = sw - int(Cm(4.0))
        cloud_height = sh - cloud_top - int(Cm(3.0)) if description else sh - cloud_top - int(Cm(1.0))

        # Arrange keywords in rows
        current_x = int(cloud_left)
        current_y = cloud_top
        row_height = 0
        padding = Cm(0.3)

        for kw in keywords:
            kw_text = kw.get("text", "")
            kw_size = kw.get("size", 24)
            kw_color_hex = kw.get("color", "#A50034")

            try:
                kr, kg, kb = int(kw_color_hex[1:3], 16), int(kw_color_hex[3:5], 16), int(kw_color_hex[5:7], 16)
                kw_color = RGBColor(kr, kg, kb)
            except (ValueError, IndexError):
                kw_color = self.colors.RED

            # Estimate box size based on text length and font size
            char_width = int(Pt(kw_size)) * 0.7
            box_w = int(len(kw_text) * char_width + int(Cm(1.0)))
            box_h = int(Pt(kw_size)) * 2 + int(Cm(0.3))

            # Wrap to next row if needed
            if current_x + box_w > int(cloud_left) + cloud_width:
                current_x = int(cloud_left)
                current_y += row_height + int(padding)
                row_height = 0

            row_height = max(row_height, box_h)

            # Background box with slight opacity (rounded)
            bg_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                current_x, current_y, box_w, box_h
            )
            bg_box.fill.solid()
            bg_box.fill.fore_color.rgb = self.colors.WHITE
            bg_box.line.color.rgb = self.colors.BORDER_GRAY
            bg_box.line.width = Pt(0.5)

            # Keyword text
            self._add_textbox(
                slide, current_x, current_y,
                box_w, box_h,
                text=kw_text, size=Pt(kw_size), bold=True,
                color=kw_color,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

            current_x += box_w + int(padding)

        # Description below
        if description:
            desc_top = sh - int(Cm(2.5))
            self._add_textbox(
                slide, Cm(2.0), desc_top,
                sw - int(Cm(4.0)), Cm(2.0),
                text=description, size=LGTypography.BODY,
                color=self.colors.DARK_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        return slide

    # ─────────────────────────────────────────
    # Swimlane Diagram
    # ─────────────────────────────────────────

    def add_swimlane(self, title, section="", subtitle="",
                     lanes=None, steps=None, connections=None, chapter=""):
        """
        Add a swimlane process diagram — horizontal lanes for each
        role/department/system with process steps placed in columns.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            lanes: List of lane names (roles/departments)
            steps: List of step dicts (see original docstring)
            connections: List of (from_step_index, to_step_index) tuples
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_start = int(y_pos)

        lanes = lanes or []
        steps = steps or []
        connections = connections or []
        if not lanes or not steps:
            return slide

        num_lanes = len(lanes)
        num_cols = max((s.get("col", 0) for s in steps), default=0) + 1

        # Layout dimensions
        label_width = Cm(3.0)
        grid_left = int(Cm(1.5)) + int(label_width) + int(Cm(0.2))
        grid_right_margin = Cm(1.0)
        grid_width = sw - grid_left - int(grid_right_margin)
        available_height = sh - y_start - int(Cm(0.5))
        lane_height = available_height // num_lanes
        col_width = grid_width // num_cols

        # Step box dimensions
        step_w = int(col_width * 0.75)
        step_h = int(lane_height * 0.55)

        # Distinct lane colors (alternating light backgrounds)
        lane_bg_colors = [
            RGBColor(0xF2, 0xF2, 0xF2),  # Light Gray
            RGBColor(0xFF, 0xFF, 0xFF),   # White
        ]

        # Lane label accent colors
        lane_accent_colors = [
            self.colors.RED,
            self.colors.CHARCOAL,
            RGBColor(0x15, 0x65, 0xC0),  # Blue
            RGBColor(0x2E, 0x7D, 0x32),  # Green
            RGBColor(0xD4, 0x76, 0x0A),  # Orange
            RGBColor(0x7B, 0x1F, 0xA2),  # Purple
        ]

        # Draw lanes (horizontal bands)
        for i, lane_name in enumerate(lanes):
            ly = y_start + lane_height * i

            # Lane background band
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                grid_left, ly, grid_width, lane_height
            )
            band.fill.solid()
            band.fill.fore_color.rgb = lane_bg_colors[i % 2]
            band.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
            band.line.width = Pt(0.5)

            # Lane label (left side, rounded)
            accent_color = lane_accent_colors[i % len(lane_accent_colors)]
            label_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(Cm(1.5)), ly, int(label_width), lane_height
            )
            label_box.fill.solid()
            label_box.fill.fore_color.rgb = accent_color
            label_box.line.fill.background()

            tf = label_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Cm(0.2)
            tf.margin_right = Cm(0.2)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = lane_name
            self._set_font(run, size=Pt(11), bold=True, color=self.colors.WHITE)

            # Thin horizontal divider at bottom of lane
            if i < num_lanes - 1:
                div_y = ly + lane_height
                div = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    grid_left, div_y, grid_width, Pt(1)
                )
                div.fill.solid()
                div.fill.fore_color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
                div.line.fill.background()

        # Place step boxes and record centers for connections
        step_centers = {}  # step_index -> (cx, cy)

        for idx, step in enumerate(steps):
            lane_idx = step.get("lane", 0)
            col_idx = step.get("col", 0)
            text = step.get("text", "")
            color_hex = step.get("color", None)

            # Calculate position (centered in cell)
            cx = grid_left + col_width * col_idx + col_width // 2
            cy = y_start + lane_height * lane_idx + lane_height // 2
            bx = cx - step_w // 2
            by = cy - step_h // 2

            # Step box color
            if color_hex:
                box_color = RGBColor(
                    int(color_hex[1:3], 16),
                    int(color_hex[3:5], 16),
                    int(color_hex[5:7], 16)
                )
            else:
                box_color = self.colors.CHARCOAL

            # Step box (rounded)
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bx, by, step_w, step_h
            )
            box.fill.solid()
            box.fill.fore_color.rgb = box_color
            box.line.fill.background()

            # Adjust corner radius
            try:
                box.adjustments[0] = 0.15
            except Exception:
                pass

            # Step text
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Cm(0.15)
            tf.margin_right = Cm(0.15)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            self._set_font(run, size=Pt(10), bold=True, color=self.colors.WHITE)

            step_centers[idx] = (cx, cy)

            self._add_shadow(box)

        # Draw connection arrows
        for from_idx, to_idx in connections:
            if from_idx not in step_centers or to_idx not in step_centers:
                continue

            fx, fy = step_centers[from_idx]
            tx, ty = step_centers[to_idx]

            from_step = steps[from_idx]
            to_step = steps[to_idx]
            from_col = from_step.get("col", 0)
            to_col = to_step.get("col", 0)
            from_lane = from_step.get("lane", 0)
            to_lane = to_step.get("lane", 0)

            if from_col == to_col:
                # Vertical arrow (same column, different lane)
                if fy < ty:
                    ay = fy + step_h // 2
                    ah = ty - step_h // 2 - ay
                else:
                    ay = ty + step_h // 2
                    ah = fy - step_h // 2 - ay

                if ah > 0:
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.DOWN_ARROW if fy < ty else MSO_SHAPE.UP_ARROW,
                        fx - int(Cm(0.2)), ay,
                        int(Cm(0.4)), ah
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = self.colors.CHARCOAL
                    arrow.line.fill.background()
            else:
                # Horizontal or diagonal → use right arrow
                ax = fx + step_w // 2
                aw = tx - step_w // 2 - ax

                if from_lane == to_lane and aw > 0:
                    # Same lane: horizontal arrow
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        ax, fy - int(Cm(0.25)),
                        aw, int(Cm(0.5))
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    arrow.line.fill.background()
                elif aw > 0:
                    # Different lane: L-shaped connector (horizontal then vertical)
                    mid_x = ax + aw // 2

                    # Horizontal segment
                    h_seg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        ax, fy - Pt(2),
                        mid_x - ax + int(Cm(0.1)), Pt(4)
                    )
                    h_seg.fill.solid()
                    h_seg.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    h_seg.line.fill.background()

                    # Vertical segment
                    v_top = min(fy, ty)
                    v_bottom = max(fy, ty)
                    v_seg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        mid_x, v_top,
                        Pt(4), v_bottom - v_top
                    )
                    v_seg.fill.solid()
                    v_seg.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    v_seg.line.fill.background()

                    # Final horizontal to target
                    h_seg2 = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        mid_x, ty - Pt(2),
                        tx - step_w // 2 - mid_x, Pt(4)
                    )
                    h_seg2.fill.solid()
                    h_seg2.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    h_seg2.line.fill.background()

                    # Arrow head at target
                    arr_head = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        tx - step_w // 2 - int(Cm(0.5)), ty - int(Cm(0.25)),
                        int(Cm(0.5)), int(Cm(0.5))
                    )
                    arr_head.fill.solid()
                    arr_head.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
                    arr_head.line.fill.background()

        return slide

    # ─────────────────────────────────────────
    # Chart / Image Slides
    # ─────────────────────────────────────────

    def add_chart_slide(self, title, section="", subtitle="",
                        chart_path=None, caption="", chapter=""):
        """
        Add a slide with a chart image (from matplotlib save or any image file).

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            chart_path: Path to chart image file (PNG, JPG, etc.)
            caption: Optional caption text below the chart
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        if chart_path and os.path.exists(chart_path):
            # Calculate image dimensions to fill content area
            img_top = y_pos + int(Cm(0.3))
            img_left = Cm(2.0)
            img_max_width = sw - int(Cm(4.0))
            img_max_height = sh - img_top - int(Cm(1.5))
            if caption:
                img_max_height -= int(Cm(1.0))

            # Add image centered
            pic = slide.shapes.add_picture(
                chart_path,
                img_left, img_top,
                width=img_max_width
            )

            # Scale to fit height if needed
            if pic.height > img_max_height:
                ratio = img_max_height / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = img_max_height

            # Center horizontally
            pic.left = (sw - pic.width) // 2

        if caption:
            cap_top = sh - int(Cm(1.5))
            self._add_textbox(
                slide, Cm(2.0), cap_top,
                sw - int(Cm(4.0)), Cm(1.0),
                text=caption, size=LGTypography.CAPTION,
                color=self.colors.MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER
            )

        return slide

    def add_image_slide(self, title, section="", subtitle="", images=None,
                        chapter=""):
        """
        Add a slide with one or more images in various layouts.

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            images: List of image dicts (see original docstring)
            chapter: Chapter name shown top-left in red
        """
        slide = self._get_blank_slide()
        sw = int(LGDimensions.SLIDE_WIDTH)
        sh = int(LGDimensions.SLIDE_HEIGHT)

        # L-Style chrome
        y_pos = self._add_content_chrome(slide, title=title, section=section,
                                         subtitle=subtitle, chapter=chapter)
        y_pos = int(y_pos)

        images = images or []
        if not images:
            return slide

        content_top = y_pos + int(Cm(0.3))
        content_height = sh - content_top - int(Cm(1.0))
        content_width = sw - int(Cm(3.0))
        margin_left = Cm(1.5)
        num_images = len(images)

        if num_images == 1:
            # Single image centered
            img = images[0]
            path = img.get("path", "")
            caption = img.get("caption", "")
            img_w = Cm(img.get("width", 20)) if img.get("width") else content_width

            if os.path.exists(path):
                pic = slide.shapes.add_picture(
                    path, int(margin_left), content_top, width=int(img_w)
                )
                # Scale to fit
                if pic.height > content_height - int(Cm(1.0)):
                    ratio = (content_height - int(Cm(1.0))) / pic.height
                    pic.width = int(pic.width * ratio)
                    pic.height = int(content_height - int(Cm(1.0)))
                # Center
                pic.left = (sw - pic.width) // 2

                if caption:
                    self._add_textbox(
                        slide, int(margin_left),
                        content_top + pic.height + int(Cm(0.2)),
                        content_width, Cm(0.8),
                        text=caption, size=LGTypography.CAPTION,
                        color=self.colors.MEDIUM_GRAY,
                        alignment=PP_ALIGN.CENTER
                    )

        elif num_images == 2:
            # Side by side
            gap = Cm(0.5)
            each_w = (content_width - int(gap)) // 2

            for i, img in enumerate(images):
                path = img.get("path", "")
                caption = img.get("caption", "")
                ix = int(margin_left) + i * (each_w + int(gap))

                if os.path.exists(path):
                    pic = slide.shapes.add_picture(
                        path, ix, content_top, width=each_w
                    )
                    if pic.height > content_height - int(Cm(1.5)):
                        ratio = (content_height - int(Cm(1.5))) / pic.height
                        pic.width = int(pic.width * ratio)
                        pic.height = int(content_height - int(Cm(1.5)))

                    if caption:
                        self._add_textbox(
                            slide, ix,
                            content_top + pic.height + int(Cm(0.2)),
                            each_w, Cm(0.8),
                            text=caption, size=LGTypography.CAPTION,
                            color=self.colors.MEDIUM_GRAY,
                            alignment=PP_ALIGN.CENTER
                        )

        else:
            # Grid layout (2 columns)
            gap = Cm(0.3)
            cols = 2
            rows_count = (num_images + cols - 1) // cols
            each_w = (content_width - int(gap) * (cols - 1)) // cols
            each_h = (content_height - int(gap) * (rows_count - 1)) // rows_count

            for i, img in enumerate(images):
                path = img.get("path", "")
                caption = img.get("caption", "")
                col = i % cols
                row = i // cols
                ix = int(margin_left) + col * (each_w + int(gap))
                iy = content_top + row * (each_h + int(gap))

                if os.path.exists(path):
                    pic = slide.shapes.add_picture(
                        path, ix, iy, width=each_w
                    )
                    if pic.height > each_h - int(Cm(1.0)):
                        ratio = (each_h - int(Cm(1.0))) / pic.height
                        pic.width = int(pic.width * ratio)
                        pic.height = int(each_h - int(Cm(1.0)))

                    if caption:
                        self._add_textbox(
                            slide, ix,
                            iy + pic.height + int(Cm(0.1)),
                            each_w, Cm(0.6),
                            text=caption, size=Pt(8),
                            color=self.colors.MEDIUM_GRAY,
                            alignment=PP_ALIGN.CENTER
                        )

        return slide

    def add_matplotlib_chart(self, title, section="", subtitle="",
                             fig=None, caption="", chapter=""):
        """
        Add a matplotlib figure directly to a slide (saves to temp file internally).

        Args:
            title: Slide title
            section: Section name for top-right
            subtitle: Optional subtitle
            fig: matplotlib Figure object
            caption: Optional caption text below the chart
            chapter: Chapter name shown top-left in red
        """
        if fig is None:
            return self.add_chart_slide(title, section, subtitle, None, caption,
                                        chapter=chapter)

        import tempfile
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        tmp_path = tmp.name
        tmp.close()

        try:
            fig.savefig(tmp_path, dpi=150, bbox_inches='tight',
                        facecolor='white', edgecolor='none')
            slide = self.add_chart_slide(title, section, subtitle,
                                         tmp_path, caption, chapter=chapter)
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

        return slide

    def save(self, filename):
        """
        Save the presentation to a file.

        Args:
            filename: Output file path.
        """
        self.prs.save(filename)
        return filename
